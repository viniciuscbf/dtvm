# -*- coding: utf-8 -*-
"""
Argus — Deck de PARCERIA para bancos.
Proposta: administradora fiduciária LOW-COST para FUNDOS PEQUENOS, operada por uma
equipe (tecnologia + operação) SOB A LICENÇA de um banco que já é instituição
autorizada. Receita começa 50/50 (banco/Argus) com opção de recompra → 15/85. Números vêm de docs auditados do projeto ou de
fontes públicas citadas — nenhum dado hipotético inventado.

Sistema de design "editorial + cores da marca" (roxo #6a50ac / azul #4a6dc0 /
ciano #1fc0ef): display serifado, fios de 1px, números tabulares, a marca em malha
da Argus (favicon), muitos diagramas. Zero blob/gradiente/glass. Gera .pptx (16:9).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_MARK
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

EMU_IN = 914400

# ---------------------------------------------------------------- paleta (marca)
NAVY   = RGBColor(0x6a, 0x50, 0xac)   # roxo da marca — cor dominante
NAVY_DK= RGBColor(0x17, 0x12, 0x3a)   # roxo quase-preto — fundos escuros
NAVY_DP= RGBColor(0x0e, 0x0a, 0x24)   # mais escuro ainda — cartões sobre escuro
TEAL   = RGBColor(0x4a, 0x6d, 0xc0)   # azul da marca — secundária
GOLD   = RGBColor(0x1f, 0xc0, 0xef)   # ciano da marca — acento
GOLD_DK= RGBColor(0x17, 0xa2, 0xce)   # ciano escuro — bordas/gradientes
INK    = RGBColor(0x2b, 0x32, 0x32)   # grafite — texto
MUTED  = RGBColor(0x5b, 0x6b, 0x7f)
FAINT  = RGBColor(0x90, 0x9d, 0xac)
LINE   = RGBColor(0xcf, 0xd8, 0xe1)
LINE_DK= RGBColor(0xe6, 0xed, 0xf3)
PAPER  = RGBColor(0xff, 0xff, 0xff)
SOFT   = RGBColor(0xf4, 0xf5, 0xfb)
CARD   = RGBColor(0xfb, 0xfc, 0xfd)
D_TXT  = RGBColor(0xf3, 0xf6, 0xfa)
D_MUT  = RGBColor(0xa8, 0xb6, 0xc6)
D_FAINT= RGBColor(0x6d, 0x80, 0x96)
D_LINE = RGBColor(0x2e, 0x26, 0x52)
POS    = RGBColor(0x2f, 0x9e, 0x7a)
NEG    = RGBColor(0xb5, 0x4a, 0x3a)

SERIF = "Georgia"
SANS  = "Calibri"
MONO  = "Consolas"

PW, PH = 13.333, 7.5
MARGIN = 0.92

prs = Presentation()
prs.slide_width  = Emu(int(PW * EMU_IN))
prs.slide_height = Emu(int(PH * EMU_IN))
BLANK = prs.slide_layouts[6]

# =================================================================== helpers base
def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    if bg is not None:
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = bg
    return s

def _set_run(run, text, size, font, color, bold, italic, char):
    run.text = text
    f = run.font
    f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color
    if char is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(int(char * 100)))

def txt(s, x, y, w, h, text, size, *, font=SANS, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, wrap=True, char=None, space_after=0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing is not None:
            p.line_spacing = spacing
        p.space_before = Pt(0); p.space_after = Pt(space_after)
        _set_run(p.add_run(), line, size, font, color, bold, italic, char)
    return tb

def rich(s, x, y, w, h, paras, *, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, (runs, opts) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get('align', PP_ALIGN.LEFT)
        p.line_spacing = opts.get('line', 1.0)
        p.space_before = Pt(opts.get('space_before', 0)); p.space_after = Pt(opts.get('space_after', 0))
        for (t, o) in runs:
            _set_run(p.add_run(), t, o.get('size', 14), o.get('font', SANS),
                     o.get('color', INK), o.get('bold', False), o.get('italic', False), o.get('char', None))
    return tb

def _soft_shadow(shp):
    spPr = shp._element.spPr
    effLst = spPr.find(qn('a:effectLst'))
    if effLst is None:
        effLst = spPr.makeelement(qn('a:effectLst'), {}); spPr.append(effLst)
    sdw = effLst.makeelement(qn('a:outerShdw'),
        {'blurRad': '90000', 'dist': '28000', 'dir': '5400000', 'rotWithShape': '0'})
    clr = sdw.makeelement(qn('a:srgbClr'), {'val': '1B2A3D'})
    alpha = clr.makeelement(qn('a:alpha'), {'val': '18000'}); clr.append(alpha); sdw.append(clr); effLst.append(sdw)

def rect(s, x, y, w, h, *, fill=None, line=None, line_w=0.75, radius=None, shadow=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if shadow: _soft_shadow(shp)
    return shp

def hline(s, x, y, w, color=LINE, weight=0.75, dash=None):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    if dash:
        d = ln.line._get_or_add_ln(); d.append(d.makeelement(qn('a:prstDash'), {'val': dash}))
    ln.shadow.inherit = False
    return ln

def vline(s, x, y, h, color=LINE, weight=0.75):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x), Inches(y + h))
    ln.line.color.rgb = color; ln.line.width = Pt(weight); ln.shadow.inherit = False
    return ln

def oval(s, cx, cy, r, *, fill=None, line=None, line_w=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def arrow(s, x1, y1, x2, y2, *, color=NAVY, weight=1.4):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color; ln.line.width = Pt(weight); ln.shadow.inherit = False
    lnEl = ln.line._get_or_add_ln()
    # tailEnd = ponta no FIM da linha (x2) → seta aponta no sentido do fluxo (esq→dir)
    lnEl.append(lnEl.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return ln

# ---------------------------------------------------------------- marca: a malha Argus
MARK = r"C:\Users\vinic\OneDrive\Desktop\dtvm\brand\favicon_argus.png"

def draw_eye(s, cx, cy, height, **_kw):
    # marca oficial (malha em degradê roxo→ciano, fundo transparente), centrada em (cx, cy).
    # favicon é quadrado, então largura = altura. Ignora kwargs de cor legados (**_kw).
    s.shapes.add_picture(MARK, Inches(cx - height / 2.0), Inches(cy - height / 2.0), height=Inches(height))

# ---------------------------------------------------------------- cromo de página
PGN = {'n': 0}
def chrome(s, section, *, dark=False, page=True):
    ink = D_TXT if dark else INK
    mut = D_MUT if dark else MUTED
    ln  = D_LINE if dark else LINE
    draw_eye(s, MARGIN + 0.13, 0.62, 0.30, outline=(D_TXT if dark else NAVY), ring=TEAL, pupil=GOLD, highlight=not dark)
    txt(s, MARGIN + 0.34, 0.47, 3.0, 0.32, "ARGUS", 12.5, font=SERIF, color=ink, bold=True, char=0.5)
    txt(s, MARGIN + 0.34, 0.47, 6.5, 0.32, section.upper(), 8.5, font=MONO, color=mut, char=1.6, align=PP_ALIGN.RIGHT)
    hline(s, MARGIN, 0.92, PW - 2 * MARGIN, color=ln, weight=0.75)
    if page:
        PGN['n'] += 1
        txt(s, PW - MARGIN - 1.5, PH - 0.52, 1.5, 0.3, f"{PGN['n']:02d}", 8.5, font=MONO,
            color=(D_FAINT if dark else FAINT), align=PP_ALIGN.RIGHT, char=1.0)
        txt(s, MARGIN, PH - 0.52, 6.0, 0.3, "ARGUS · PROPOSTA DE PARCERIA · CONFIDENCIAL", 8, font=MONO,
            color=(D_FAINT if dark else FAINT), char=1.2)

def kicker(s, x, y, text, color=GOLD, w=8.0):
    hline(s, x, y + 0.09, 0.28, color=color, weight=1.6)
    txt(s, x + 0.4, y, w, 0.3, text.upper(), 9.5, font=MONO, color=color, char=2.0)

def head(s, section, kick, title, *, dark=False, tsize=29, th=0.95, kcolor=GOLD):
    chrome(s, section, dark=dark)
    kicker(s, MARGIN, 1.22, kick, color=kcolor)
    txt(s, MARGIN, 1.56, PW - 2 * MARGIN, th, title, tsize, font=SERIF, color=(D_TXT if dark else INK), spacing=1.02)
    return 1.56 + th + 0.22

def card(s, x, y, w, h, *, fill=CARD, line=LINE, line_w=0.75, radius=0.045, shadow=True):
    return rect(s, x, y, w, h, fill=fill, line=line, line_w=line_w, radius=radius, shadow=shadow)

def numeral(s, x, y, n, *, color=GOLD, size=15):
    txt(s, x, y, 0.9, 0.4, f"{n:02d}", size, font=MONO, color=color, bold=True, char=0.5)

def source(s, text, *, dark=False):
    txt(s, MARGIN, PH - 0.78, PW - 2 * MARGIN, 0.3, text, 8, font=MONO,
        color=(D_FAINT if dark else FAINT), char=0.2)

# ---------------------------------------------------------------- gráficos
def _cf(chart, color=MUTED):
    chart.font.name = MONO; chart.font.size = Pt(9); chart.font.color.rgb = color

def line_chart(s, x, y, w, h, cats, series, *, colors=None, smooth=False, markers=True,
               val_fmt='0.0', legend=False, axis_color=MUTED, grid=True, widths=None):
    cd = CategoryChartData(); cd.categories = cats
    for name, vals in series: cd.add_series(name, vals, number_format=val_fmt)
    gf = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS if markers else XL_CHART_TYPE.LINE,
                            Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart; ch.has_title = False; ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.TOP; ch.legend.include_in_layout = False
        ch.legend.font.name = MONO; ch.legend.font.size = Pt(9); ch.legend.font.color.rgb = MUTED
    _cf(ch, axis_color)
    cols = colors or [NAVY, TEAL, GOLD]
    for i, ps in enumerate(ch.series):
        c = cols[i % len(cols)]; ps.smooth = smooth
        lf = ps.format.line; lf.color.rgb = c; lf.width = Pt((widths or [2.4])[i % len(widths or [2.4])])
        try:
            m = ps.marker; m.style = 8 if markers else 2; m.size = 6
            m.format.fill.solid(); m.format.fill.fore_color.rgb = c
            m.format.line.color.rgb = PAPER; m.format.line.width = Pt(1.0)
        except Exception: pass
    ca = ch.category_axis; va = ch.value_axis
    ca.tick_labels.font.name = MONO; ca.tick_labels.font.size = Pt(9); ca.tick_labels.font.color.rgb = axis_color
    va.tick_labels.font.name = MONO; va.tick_labels.font.size = Pt(9); va.tick_labels.font.color.rgb = axis_color
    ca.format.line.color.rgb = LINE
    ca.major_tick_mark = XL_TICK_MARK.NONE; ca.minor_tick_mark = XL_TICK_MARK.NONE
    va.major_tick_mark = XL_TICK_MARK.NONE; va.minor_tick_mark = XL_TICK_MARK.NONE
    va.has_major_gridlines = grid
    if grid:
        gl = va.major_gridlines.format.line; gl.color.rgb = LINE_DK; gl.width = Pt(0.5)
    ca.has_major_gridlines = False; va.format.line.fill.background()
    return ch

def bar_chart(s, x, y, w, h, cats, series, *, colors=None, col=True, val_fmt='0', legend=False,
              show_val=False, axis_color=MUTED, gap=60, grid=True, val_color=INK, overlap=-20):
    cd = CategoryChartData(); cd.categories = cats
    for name, vals in series: cd.add_series(name, vals, number_format=val_fmt)
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED if col else XL_CHART_TYPE.BAR_CLUSTERED,
                            Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart; ch.has_title = False; ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.TOP; ch.legend.include_in_layout = False
        ch.legend.font.name = MONO; ch.legend.font.size = Pt(9); ch.legend.font.color.rgb = MUTED
    _cf(ch, axis_color)
    plot = ch.plots[0]; plot.gap_width = gap
    try: plot.overlap = overlap
    except Exception: pass
    cols = colors or [NAVY, TEAL, GOLD]
    for i, ser in enumerate(ch.series):
        ser.format.fill.solid(); ser.format.fill.fore_color.rgb = cols[i % len(cols)]; ser.format.line.fill.background()
    if show_val:
        plot.has_data_labels = True; dl = plot.data_labels
        dl.number_format = val_fmt; dl.number_format_is_linked = False
        dl.font.name = MONO; dl.font.size = Pt(9); dl.font.bold = True; dl.font.color.rgb = val_color
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
    ca = ch.category_axis; va = ch.value_axis
    ca.tick_labels.font.name = MONO; ca.tick_labels.font.size = Pt(9); ca.tick_labels.font.color.rgb = axis_color
    va.tick_labels.font.name = MONO; va.tick_labels.font.size = Pt(9); va.tick_labels.font.color.rgb = axis_color
    ca.format.line.color.rgb = LINE
    ca.major_tick_mark = XL_TICK_MARK.NONE; ca.minor_tick_mark = XL_TICK_MARK.NONE
    va.major_tick_mark = XL_TICK_MARK.NONE; va.minor_tick_mark = XL_TICK_MARK.NONE
    va.has_major_gridlines = grid
    if grid:
        gl = va.major_gridlines.format.line; gl.color.rgb = LINE_DK; gl.width = Pt(0.5)
    if not grid: va.visible = False
    return ch

def color_points(chart, colors):
    ser = chart.series[0]
    try: ser.invert_if_negative = False
    except Exception: pass
    for i, pt in enumerate(ser.points):
        pt.format.fill.solid(); pt.format.fill.fore_color.rgb = colors[i]; pt.format.line.fill.background()
        dpt = pt._element
        if dpt.find(qn('c:invertIfNegative')) is None:
            dpt.find(qn('c:idx')).addnext(dpt.makeelement(qn('c:invertIfNegative'), {'val': '0'}))

def stacked_col(s, x, y, w, h, cats, series, colors, *, val_fmt='0', legend=True, show_val=True, horiz=False):
    cd = CategoryChartData(); cd.categories = cats
    for name, vals in series: cd.add_series(name, vals, number_format=val_fmt)
    ct = XL_CHART_TYPE.BAR_STACKED if horiz else XL_CHART_TYPE.COLUMN_STACKED
    gf = s.shapes.add_chart(ct, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart; ch.has_title = False; ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.RIGHT; ch.legend.include_in_layout = False
        ch.legend.font.name = MONO; ch.legend.font.size = Pt(9); ch.legend.font.color.rgb = MUTED
    _cf(ch, MUTED); ch.plots[0].gap_width = 80
    for i, ser in enumerate(ch.series):
        ser.format.fill.solid(); ser.format.fill.fore_color.rgb = colors[i % len(colors)]
        ser.format.line.color.rgb = PAPER; ser.format.line.width = Pt(1.0)
    if show_val:
        for ser in ch.series:
            ser.has_data_labels = True; dl = ser.data_labels
            dl.number_format = val_fmt; dl.number_format_is_linked = False
            dl.font.name = MONO; dl.font.size = Pt(9); dl.font.bold = True; dl.font.color.rgb = PAPER
            dl.position = XL_LABEL_POSITION.CENTER
    ca = ch.category_axis; va = ch.value_axis
    ca.tick_labels.font.name = MONO; ca.tick_labels.font.size = Pt(10); ca.tick_labels.font.color.rgb = MUTED
    va.visible = False; ca.format.line.color.rgb = LINE
    ca.major_tick_mark = XL_TICK_MARK.NONE; va.has_major_gridlines = False
    return ch

def table(s, x, y, cols, headers, rows, *, rh=0.6, accent_col=None, hsize=9, csize=10, dsize=11.5):
    xs = []; cx = x
    for w in cols: xs.append(cx); cx += w
    total = cx - x
    for i, htxt in enumerate(headers):
        col = GOLD if i == accent_col else MUTED
        txt(s, xs[i], y, cols[i] - 0.2, 0.3, htxt, hsize, font=MONO, color=col, char=1.1, bold=(i == accent_col))
    yy = y + 0.32
    hline(s, x, yy, total, color=INK, weight=1.0); yy += 0.12
    for r in rows:
        for i, cell in enumerate(r):
            if i == 0:
                txt(s, xs[i], yy + 0.03, cols[i] - 0.2, rh, cell, dsize, font=SERIF, color=INK, bold=True, spacing=1.05)
            else:
                col = INK if i == accent_col else MUTED
                txt(s, xs[i], yy + 0.03, cols[i] - 0.2, rh, cell, csize, font=SANS, color=col, spacing=1.15)
        yy += rh
        hline(s, x, yy, total, color=LINE_DK, weight=0.75); yy += 0.05
    return yy

def pipeline(s, x, y, w, steps, *, h=0.95, gap=0.26, dark=False, highlight=None, tsize=10.5):
    """Fluxo horizontal de caixas ligadas por setas. steps=[(titulo, sub)]. highlight=set de índices."""
    highlight = highlight or set()
    n = len(steps)
    bw = (w - (n - 1) * gap) / n
    line_n = D_LINE if dark else LINE
    scol = D_MUT if dark else MUTED
    for i, (t, sub) in enumerate(steps):
        bx = x + i * (bw + gap); hi = i in highlight
        f = NAVY if hi else (NAVY_DP if dark else CARD)
        rect(s, bx, y, bw, h, fill=f, line=(GOLD if hi else line_n), line_w=(1.3 if hi else 0.75),
             radius=0.06, shadow=not dark)
        ttc = D_TXT if (hi or dark) else INK
        txt(s, bx + 0.1, y + 0.12, bw - 0.2, 0.55, t, tsize, font=SANS, color=ttc, bold=True,
            align=PP_ALIGN.CENTER, spacing=1.0, anchor=MSO_ANCHOR.TOP)
        if sub:
            txt(s, bx + 0.08, y + h - 0.34, bw - 0.16, 0.3, sub, 8, font=MONO,
                color=(GOLD if hi else scol), align=PP_ALIGN.CENTER, char=0.2)
        if i < n - 1:
            ax = bx + bw + 0.02
            arrow(s, ax, y + h / 2, ax + gap - 0.05, y + h / 2, color=(D_FAINT if dark else FAINT), weight=1.3)

# =================================================================== DECK
def s_cover():
    s = slide(NAVY_DK)
    draw_eye(s, MARGIN + 0.5, 1.45, 1.15, outline=D_TXT, ring=TEAL, pupil=GOLD, highlight=True)
    txt(s, MARGIN + 1.22, 1.08, 6, 0.9, "ARGUS", 33, font=SERIF, color=D_TXT, bold=True, char=1.0)
    txt(s, MARGIN + 1.25, 1.8, 9, 0.4, "ADMINISTRAÇÃO FIDUCIÁRIA DE FUNDOS PEQUENOS", 9.5,
        font=MONO, color=D_MUT, char=2.0)
    hline(s, MARGIN, 3.0, PW - 2 * MARGIN, color=D_LINE, weight=0.75)
    rich(s, MARGIN, 3.32, 11.2, 2.0, [
        ([("Uma nova linha de receita recorrente", {'size': 37, 'font': SERIF, 'color': D_TXT})], {'line': 1.05}),
        ([("para o seu banco.", {'size': 37, 'font': SERIF, 'color': GOLD, 'italic': True})],
         {'line': 1.05, 'space_before': 2}),
    ])
    txt(s, MARGIN, 5.35, 10.2, 1.0,
        "Vocês entram com a licença e a estrutura que já têm. Nós entramos com a tecnologia e a "
        "operação. Juntos, atendemos um mercado que os grandes decidiram ignorar — os fundos pequenos.",
        14, font=SANS, color=D_MUT, spacing=1.32)
    hline(s, MARGIN, PH - 0.82, PW - 2 * MARGIN, color=D_LINE, weight=0.75)
    txt(s, MARGIN, PH - 0.68, 8, 0.3, "PROPOSTA DE PARCERIA · CONFIDENCIAL", 8.5, font=MONO, color=D_FAINT, char=1.8)
    txt(s, PW - MARGIN - 5, PH - 0.68, 5, 0.3, "PROTÓTIPO OPERACIONAL · DADOS SIMULADOS · 2026", 8.5,
        font=MONO, color=D_FAINT, char=1.6, align=PP_ALIGN.RIGHT)

def s_problema():
    s = slide(PAPER)
    head(s, "O problema", "A lacuna de mercado", "A regulação permite o fundo de R$ 1 milhão. O custo, não.")
    txt(s, MARGIN, 2.62, 4.75, 3.2,
        "Um fundo pequeno carrega hoje R$ 50–70 mil por ano de custo fixo — gestão, administração, "
        "custódia, auditoria e taxa CVM — quase todo independente do tamanho do fundo.\n\n"
        "Num fundo de R$ 1 milhão, isso é 5–7% do patrimônio ao ano. Nenhum fundo sobrevive a isso. "
        "Por isso, na prática, só se abre um fundo a partir de ~R$ 10 milhões.", 13, font=SANS, color=INK, spacing=1.34)
    txt(s, 6.35, 2.6, 6.2, 0.3, "CUSTO FIXO ANUAL (~R$ 60 MIL) COMO % DO PL DO FUNDO", 9,
        font=MONO, color=MUTED, char=1.0)
    ch = bar_chart(s, 6.2, 2.95, PW - MARGIN - 6.2, 3.35,
                   ["R$ 1 mi", "R$ 2 mi", "R$ 5 mi", "R$ 10 mi", "R$ 20 mi"],
                   [("% do PL", [6.0, 3.0, 1.2, 0.6, 0.3])],
                   colors=[NAVY], col=True, show_val=True, val_fmt='0.0"%"', grid=False, gap=55, val_color=NAVY)
    color_points(ch, [NEG, NEG, NAVY, NAVY, NAVY])
    source(s, "Custo fixo de um fundo pequeno: R$ 50–70 mil/ano (InfoMoney, 2023). % calculado sobre R$ 60 mil.")

def s_lacuna():
    s = slide(PAPER)
    head(s, "O problema", "Por que a lacuna existe", "O custo de operar um fundo é fixo — não cabe no pequeno.")
    pipeline(s, MARGIN, 2.75, PW - 2 * MARGIN,
             [("Custo por fundo\né majoritariamente FIXO", "cota · conciliação · reporte"),
              ("Não cai\ncom o tamanho do fundo", "R$ mil por fundo"),
              ("Administradora cobra\nmínimo mensal alto", "milhares/mês"),
              ("Fundo pequeno\nfica inviável", "5–7% do PL"),
              ("Os grandes\nignoram o segmento", "foco em PL alto")],
             h=1.35, highlight={4}, tsize=10.5)
    hline(s, MARGIN, 4.75, PW - 2 * MARGIN, color=LINE)
    rich(s, MARGIN, 5.0, PW - 2 * MARGIN, 1.4, [
        ([("A ineficiência não é maldade — é estrutura de custo. ", {'size': 14, 'font': SANS, 'color': INK, 'bold': True}),
          ("Calcular cota, conciliar e reportar custa quase o mesmo para um fundo de R$ 1 milhão e para um "
           "de R$ 1 bilhão. Como o administrador tradicional faz isso com gente e sistemas antigos, ele "
           "precisa de mínimos altos — e o fundo pequeno simplesmente não cabe na conta dele.",
           {'size': 14, 'font': SANS, 'color': MUTED})], {'line': 1.34})])
    source(s, "Diagnóstico: plano_dtvm_fundos_pequenos.md (§2).")

def s_oportunidade():
    s = slide(PAPER)
    head(s, "A oportunidade", "Quem fica de fora", "Existe demanda — represada pelo custo.")
    cards = [
        ("01", "Gestores certificados", "Profissionais que começariam com R$ 2–5 milhões e uma estratégia própria — mas não encontram administrador que os atenda por preço viável."),
        ("02", "Assessorias de investimento", "27.515 assessores e 1.418 escritórios que poderiam ter o fundo próprio como produto — hoje sem uma via barata para criá-lo."),
        ("03", "Estratégias de nicho", "Casas pequenas: 51% das gestoras têm menos de R$ 400 milhões. Só 1.277 dos ~31 mil fundos passam de R$ 1 bilhão — a cauda é imensa."),
    ]
    x0 = MARGIN; cw = (PW - 2 * MARGIN - 2 * 0.45) / 3
    for i, (n, t, d) in enumerate(cards):
        x = x0 + i * (cw + 0.45)
        card(s, x, 2.75, cw, 3.05, radius=0.04)
        numeral(s, x + 0.32, 3.02, int(n), size=15)
        hline(s, x + 0.34, 3.5, 0.32, color=GOLD, weight=1.6)
        txt(s, x + 0.32, 3.62, cw - 0.64, 0.5, t, 15, font=SERIF, color=INK, bold=True)
        txt(s, x + 0.32, 4.18, cw - 0.64, 1.5, d, 11, font=SANS, color=MUTED, spacing=1.28)
    rect(s, MARGIN, 6.05, PW - 2 * MARGIN, 0.62, fill=SOFT, line=LINE, radius=0.06, shadow=False)
    txt(s, MARGIN + 0.25, 6.17, PW - 2 * MARGIN - 0.5, 0.45,
        "E os concorrentes de infraestrutura (Kanastra, QI Tech, Vórtx, BRLTrust, Singulare) foram todos "
        "para FIDC e crédito estruturado de ticket alto. O fundo pequeno genérico segue desatendido.",
        10.5, font=SANS, color=INK, italic=True, spacing=1.15)
    source(s, "ANCORD (assessores, dez/2025); ANBIMA (976 gestoras; 51% <R$400mi); imprensa (foco das fintechs).")

def s_tese():
    s = slide(NAVY_DK)
    head(s, "A tese", "A virada", "Construir para escala desde o primeiro fundo.", dark=True)
    txt(s, MARGIN, 2.5, PW - 2 * MARGIN, 0.6,
        "Se a infraestrutura for desenhada para rodar centenas de fundos, o custo do fundo nº 500 é quase "
        "o do fundo nº 5. Aí o segmento que ninguém quer vira lucrativo.", 14, font=SANS, color=D_MUT, spacing=1.3)
    # dois caminhos
    y = 3.7
    card(s, MARGIN, y, 5.35, 2.5, fill=NAVY_DP, line=D_LINE, radius=0.05, shadow=False)
    txt(s, MARGIN + 0.3, y + 0.22, 4.8, 0.3, "HOJE — ESTRUTURA TRADICIONAL", 9, font=MONO, color=NEG, char=1.2)
    for i, t in enumerate(["Custo fixo alto por fundo", "Mínimo mensal caro", "Fundo pequeno inviável"]):
        yy = y + 0.7 + i * 0.52
        txt(s, MARGIN + 0.3, yy, 0.3, 0.3, "—", 12, font=SANS, color=D_FAINT)
        txt(s, MARGIN + 0.6, yy, 4.5, 0.4, t, 12, font=SANS, color=D_TXT, spacing=1.05)
    arrow(s, MARGIN + 5.55, y + 1.25, MARGIN + 6.35, y + 1.25, color=GOLD, weight=2.2)
    card(s, MARGIN + 6.55, y, 5.35, 2.5, fill=NAVY, line=GOLD, radius=0.05, shadow=False)
    txt(s, MARGIN + 6.85, y + 0.22, 4.8, 0.3, "A PROPOSTA — INFRA PARA ESCALA", 9, font=MONO, color=GOLD, char=1.2)
    for i, t in enumerate(["Custo marginal por fundo ~zero", "Taxa enxuta: 0,08% + R$100/mês", "Fundo pequeno viável e lucrativo"]):
        yy = y + 0.7 + i * 0.52
        txt(s, MARGIN + 6.85, yy, 0.3, 0.3, "—", 12, font=SANS, color=GOLD)
        txt(s, MARGIN + 7.15, yy, 4.5, 0.4, t, 12, font=SANS, color=D_TXT, spacing=1.05)
    source(s, "plano_dtvm_fundos_pequenos.md (§2–3). Taxa 0,08% a.a. + piso R$ 100/mês.", dark=True)

def s_solucao():
    s = slide(PAPER)
    head(s, "A solução", "A plataforma", "Uma administradora onde a máquina faz o ciclo inteiro.")
    pipeline(s, MARGIN, 2.7, PW - 2 * MARGIN,
             [("Batch D-1\nprecifica ativos", "por curvas"),
              ("Concilia com\no custodiante", "caixa e ativos"),
              ("Calcula\na cota", "automático"),
              ("Prévia ao\ngestor", "D-1"),
              ("Gestor\naprova", "dupla checagem"),
              ("Publica\na cota", "D-0"),
              ("Informe à\nCVM", "1 dia útil")],
             h=1.3, highlight={4}, tsize=10, gap=0.2)
    rect(s, MARGIN, 4.55, PW - 2 * MARGIN, 0.66, fill=SOFT, line=GOLD, radius=0.05, shadow=False)
    txt(s, MARGIN + 0.25, 4.68, PW - 2 * MARGIN - 0.5, 0.45,
        "Nada é publicado sem a aprovação formal do gestor. Cada correção vira lançamento com trilha. "
        "É a dupla checagem que produz prova de diligência — a blindagem do banco.", 11,
        font=SANS, color=INK, italic=True, spacing=1.15)
    extras = [
        ("Custódia integrada", "o gestor boleta, a liquidação é entrega-contra-pagamento e a posição reflete no dia."),
        ("Reporte automático", "informe diário à CVM em 1 dia útil; mensais em 10 — sem trabalho manual."),
        ("Camada de IA", "leitura de regulamentos e monitoramento de fraude em 100% dos fundos, todo dia."),
    ]
    x0 = MARGIN; cw = (PW - 2 * MARGIN - 2 * 0.4) / 3
    for i, (t, d) in enumerate(extras):
        x = x0 + i * (cw + 0.4)
        oval(s, x + 0.06, 5.65, 0.05, fill=GOLD)
        txt(s, x + 0.24, 5.5, cw - 0.24, 0.35, t, 12.5, font=SERIF, color=INK, bold=True)
        txt(s, x + 0.24, 5.85, cw - 0.24, 0.8, d, 10.5, font=SANS, color=MUTED, spacing=1.22)
    source(s, "Ciclo validado no protótipo; espelha o fluxo da Res. CVM 175 (informe diário em 1 d.u.).")

def s_deploy():
    s = slide(PAPER)
    head(s, "A solução", "Diferencial", "O fundo vira um instrumento ágil — de abrir e de encerrar.")
    pipeline(s, MARGIN, 2.7, PW - 2 * MARGIN,
             [("Solicitação\ndo cliente", None),
              ("Constituição\npor templates", "regulamento · políticas"),
              ("Registro +\nintegrações", "automáticas"),
              ("Operação", "cota · enquadramento"),
              ("Encerramento\npadronizado", "baixo atrito")],
             h=1.25, highlight={1}, tsize=10.5)
    cols = [
        ("Menor tempo até operar", "colocar um fundo de pé em prazo muito menor que o processo artesanal de hoje."),
        ("Baixo custo de saída", "encerrar um fundo que não vingou deixa de ser um processo caro e travado."),
        ("Experimentação viável", "criar e encerrar com agilidade reduz o risco de tentar — combina com fundos pequenos."),
    ]
    x0 = MARGIN; cw = (PW - 2 * MARGIN - 2 * 0.45) / 3
    for i, (t, d) in enumerate(cols):
        x = x0 + i * (cw + 0.45)
        card(s, x, 4.35, cw, 1.75, radius=0.04)
        txt(s, x + 0.28, 4.55, cw - 0.56, 0.4, t, 13.5, font=SERIF, color=INK, bold=True)
        txt(s, x + 0.28, 5.02, cw - 0.56, 0.9, d, 10.5, font=SANS, color=MUTED, spacing=1.24)
    txt(s, MARGIN, 6.35, PW - 2 * MARGIN, 0.5,
        "Honesto: “ágil” é o mais rápido que a regulação permite. O encerramento tem etapas legais "
        "(liquidação, pagamento, baixa na CVM) que não se comprimem à vontade — mas sem gargalo interno.",
        9.5, font=SANS, color=MUTED, italic=True, spacing=1.15)
    source(s, "plano_dtvm_fundos_pequenos.md (§4.3).")

def s_captacao():
    s = slide(PAPER)
    head(s, "A captação", "De onde vêm os fundos", "Quatro portas de entrada — e uma escada.")
    txt(s, MARGIN, 2.58, PW - 2 * MARGIN, 0.28, "AGORA · FUNDOS PEQUENOS (R$ 1–10 MILHÕES)", 9,
        font=MONO, color=NAVY, char=1.6, bold=True)
    seg = [
        ("01", "Clubes de investimento",
         "Batem no teto de 50 participantes e querem crescer. Viram fundo para captar do público — e, por subclasse, o sócio-fundador paga taxa menor que o público novo.",
         "subclasse: sócio paga menos"),
        ("02", "Assessores de investimento",
         "Milhares de escritórios que hoje alocam cliente a cliente passam a distribuir um veículo dedicado — com gestor habilitado — para concentrar o capital dos clientes num fundo só.",
         "distribuidor · gestor habilitado"),
        ("03", "Gestores",
         "Fatiam uma estratégia em vários veículos sob medida — por benchmark e por perfil de cliente — e multiplicam a captação. É o multiplicador do próximo slide.",
         "1 fundo → N veículos"),
    ]
    cw = (PW - 2 * MARGIN - 2 * 0.45) / 3
    for i, (n, t, d, tag) in enumerate(seg):
        x = MARGIN + i * (cw + 0.45)
        card(s, x, 2.9, cw, 2.45, radius=0.04)
        numeral(s, x + 0.3, 3.12, int(n), size=14)
        txt(s, x + 0.3, 3.56, cw - 0.6, 0.35, t, 14.5, font=SERIF, color=INK, bold=True)
        txt(s, x + 0.3, 3.98, cw - 0.6, 1.05, d, 10, font=SANS, color=MUTED, spacing=1.22)
        hline(s, x + 0.3, 5.0, cw - 0.6, color=LINE_DK, weight=0.75)
        oval(s, x + 0.36, 5.14, 0.04, fill=GOLD)
        txt(s, x + 0.5, 5.04, cw - 0.8, 0.3, tag, 9.5, font=MONO, color=NAVY, char=0.3)
    rect(s, MARGIN, 5.6, PW - 2 * MARGIN, 1.02, fill=SOFT, line=NAVY, radius=0.05, shadow=False)
    txt(s, MARGIN + 0.32, 5.73, 3.4, 0.3, "2º PASSO", 11, font=MONO, color=NAVY, char=1.8, bold=True)
    txt(s, MARGIN + 0.32, 6.0, 4.6, 0.5, "Fundos médios · R$ 50–800 mi", 15, font=SERIF, color=INK, bold=True)
    txt(s, MARGIN + 5.1, 5.74, PW - 2 * MARGIN - 5.35, 0.85,
        "Com a operação calibrada, subimos de ticket. A mesma máquina que viabiliza um fundo de R$ 1 mi "
        "barateia enormemente um de R$ 200 mi. Acima de ~R$ 800 mi, o fundo já negocia taxa com os grandes — não é alvo.",
        10, font=SANS, color=INK, spacing=1.18)
    source(s, "plano_dtvm_fundos_pequenos.md (§6.2). Captação por segmento — três portas agora, os médios como 2º passo.")

def s_subclasses():
    s = slide(PAPER)
    head(s, "A captação", "O multiplicador", "Um fundo, vários veículos sob medida.")
    txt(s, MARGIN, 2.58, PW - 2 * MARGIN, 0.5,
        "O gestor não precisa de um fundo novo para cada público. Sob uma única casca, a plataforma cria "
        "quantos veículos ele quiser — um para cada perfil — a custo marginal quase zero.",
        12, font=SANS, color=MUTED, spacing=1.25)
    fb_x, fb_y, fb_w, fb_h = MARGIN, 3.55, 3.35, 2.15
    card(s, fb_x, fb_y, fb_w, fb_h, fill=NAVY, line=NAVY, radius=0.05)
    txt(s, fb_x + 0.3, fb_y + 0.26, fb_w - 0.6, 0.3, "1 FUNDO", 11, font=MONO, color=GOLD, char=1.4, bold=True)
    txt(s, fb_x + 0.3, fb_y + 0.58, fb_w - 0.6, 0.55, "A casca", 20, font=SERIF, color=D_TXT, bold=True)
    for i, t in enumerate(["1 CNPJ · 1 regulamento", "1 auditoria · 1 contabilidade", "1 registro CVM/ANBIMA"]):
        txt(s, fb_x + 0.3, fb_y + 1.24 + i * 0.28, fb_w - 0.6, 0.28, t, 9.5, font=SANS, color=D_MUT)
    chips = [
        ("SUBCLASSE", "Sócios do clube", "taxa de gestão menor"),
        ("SUBCLASSE", "Público geral", "taxa cheia"),
        ("CLASSE", "Perfil conservador", "benchmark CDI"),
        ("CLASSE", "Perfil arrojado", "benchmark Bovespa + S&P"),
    ]
    cx = 6.95; cwid = PW - MARGIN - cx; chh = 0.6; step = 0.72; cy0 = 3.35
    for i, (tag, t, d) in enumerate(chips):
        cy = cy0 + i * step
        arrow(s, fb_x + fb_w, fb_y + fb_h / 2, cx, cy + chh / 2, color=(GOLD if i < 2 else TEAL), weight=1.3)
        card(s, cx, cy, cwid, chh, radius=0.04)
        txt(s, cx + 0.24, cy + 0.1, 1.7, 0.28, tag, 8, font=MONO, color=(GOLD if i < 2 else TEAL), char=1.0, bold=True)
        txt(s, cx + 0.24, cy + 0.29, 3.0, 0.3, t, 12.5, font=SERIF, color=INK, bold=True)
        txt(s, cx + 2.9, cy + 0.18, cwid - 3.1, 0.35, d, 10.5, font=SANS, color=MUTED, align=PP_ALIGN.RIGHT)
    txt(s, MARGIN, 6.32, PW - 2 * MARGIN, 0.35,
        "Classe separa ATIVOS (estratégias distintas); subclasse separa PASSIVO (público, taxas, benchmark) — mesmo portfólio. "
        "Na infra tradicional, cada um seria um fundo novo e caro.",
        10.5, font=SANS, color=INK, italic=True, spacing=1.1)
    source(s, "guia_estruturas_classes_subclasses.md (Res. CVM 175). Classe = CNPJ/ativos; subclasse = só passivo.")

def s_preco():
    s = slide(NAVY_DK)
    head(s, "O produto", "O preço", "Muito mais barato — porque o custo é de software.", dark=True)
    # bloco do preço
    card(s, MARGIN, 2.7, 5.15, 3.5, fill=NAVY_DP, line=D_LINE, radius=0.05, shadow=False)
    txt(s, MARGIN + 0.35, 3.0, 4.5, 0.3, "A TAXA DA PLATAFORMA", 9, font=MONO, color=GOLD, char=1.4)
    txt(s, MARGIN + 0.35, 3.4, 4.5, 0.7, "0,08%", 44, font=SERIF, color=D_TXT, bold=True)
    txt(s, MARGIN + 0.35, 4.25, 4.5, 0.3, "ao ano sobre o patrimônio", 12, font=SANS, color=D_MUT)
    hline(s, MARGIN + 0.35, 4.75, 4.45, color=D_LINE)
    txt(s, MARGIN + 0.35, 4.9, 4.5, 0.5, "+ piso de R$ 100/mês por fundo", 15, font=SERIF, color=GOLD, bold=True)
    txt(s, MARGIN + 0.35, 5.5, 4.5, 0.6, "Só a administração é padronizada pela plataforma. Gestão e "
        "performance ficam com o gestor, no regulamento.", 10, font=SANS, color=D_MUT, spacing=1.2)
    # comparativo custo all-in de um fundo de R$1mi
    txt(s, 6.5, 2.6, 6.2, 0.3, "CUSTO TOTAL ANUAL DE UM FUNDO DE R$ 1 MILHÃO", 9, font=MONO, color=D_MUT, char=1.0)
    bar_chart(s, 6.35, 3.0, PW - MARGIN - 6.35, 3.15,
              ["Tradicional", "Com a Argus"],
              [("R$ mil/ano", [60, 13])],
              colors=[NAVY], col=False, show_val=True, val_fmt='"R$ "0" mil"', grid=False, gap=80, val_color=D_TXT, axis_color=D_MUT)
    txt(s, 6.35, 6.3, 6.2, 0.4, "Tradicional: R$ 50–70 mil/ano (InfoMoney). Argus: 1,29% de R$ 1 mi (custo "
        "total do modelo, planilha de custos).", 8.5, font=MONO, color=D_FAINT, char=0.2, spacing=1.15)

def s_barateamento():
    s = slide(PAPER)
    head(s, "O produto", "Onde o custo mora", "Onde cada real some.")
    rect(s, MARGIN, 2.6, PW - 2 * MARGIN, 0.66, fill=SOFT, line=GOLD, radius=0.05, shadow=False)
    txt(s, MARGIN + 0.28, 2.72, PW - 2 * MARGIN - 0.56, 0.45,
        "O piso de administração dos grandes é R$ 15–25 mil por mês (BRL Trust, Vórtx, Reag, Oliveira Trust). "
        "O nosso é R$ 100. A diferença não é o percentual — é o PISO.", 11.5, font=SANS, color=INK, italic=True, spacing=1.12)
    table(s, MARGIN, 3.55, [2.75, 2.15, 1.9, PW - 2 * MARGIN - 2.75 - 2.15 - 1.9],
          ["COMPONENTE (fundo R$ 1 mi)", "MERCADO", "ARGUS", "COMO BARATEAMOS"],
          [["Administração", "R$ 15–25 mil/mês", "R$ 100/mês", "software, não trabalho manual por fundo"],
           ["Custódia", "~R$ 1.500/mês", "R$ 0", "absorvida pelo banco parceiro"],
           ["Auditoria", "~R$ 11 mil/ano", "~R$ 1,5 mil/ano", "mesmo auditor em lote, carteiras padronizadas"],
           ["Taxa CVM", "~R$ 3 mil/ano", "~R$ 3 mil/ano", "igual — é tributo, não dá para baratear"],
           ["Total ex-gestão", "~R$ 50–70 mil/ano", "~R$ 6 mil/ano", "≈ 10× mais barato no fundo pequeno"]],
          rh=0.48, accent_col=2, csize=9.5, dsize=11)
    source(s, "Regulamentos reais (CVM/Fundos.NET): BRL Trust, Vórtx, Oliveira Trust, Singulare, Reag; auditoria XP. Ver benchmark_precos_concorrencia.md.")

def s_cotista():
    s = slide(PAPER)
    head(s, "O produto", "O produto vende", "E o fundo entrega ao cotista? Sim — a partir de R$ 2 milhões.")
    txt(s, MARGIN, 2.6, 4.7, 3.4,
        "Um fundo de crédito privado rendendo CDI+1, com gestão de 0,7% e os nossos custos, entrega ao "
        "cotista, líquido de tudo:\n\n"
        "•  R$ 1 mi — 97% do CDI (caso-limite)\n"
        "•  R$ 2 mi — empata com o CDI\n"
        "•  R$ 5 mi + — supera o CDI\n\n"
        "O vilão do fundo minúsculo são custos fixos que não controlamos (auditoria, taxa CVM). Existe "
        "produto vendável exatamente na faixa que ninguém atende.", 12, font=SANS, color=INK, spacing=1.3)
    txt(s, 6.35, 2.6, 6.2, 0.3, "RENTABILIDADE LÍQUIDA COMO % DO CDI, POR PL DO FUNDO", 9, font=MONO, color=MUTED, char=1.0)
    ch = line_chart(s, 6.2, 2.95, PW - MARGIN - 6.2, 3.35,
                    ["R$1mi", "R$2mi", "R$5mi", "R$10mi", "R$50mi"],
                    [("% do CDI", [97.3, 99.9, 101.2, 101.5, 101.9]), ("CDI (100%)", [100, 100, 100, 100, 100])],
                    colors=[NAVY, FAINT], smooth=False, val_fmt='0.0"%"', widths=[2.6, 1.0])
    source(s, "planilha_custos.md (§7.1, regime permanente). Premissa: CDI 10,5%, ativos CDI+1; conservador com a Selic atual (14,25%).")

def s_modelo():
    s = slide(PAPER)
    head(s, "A parceria", "O modelo", "Cada lado faz o que já sabe fazer.")
    lw = 5.3; rxx = MARGIN + 6.15; rw = PW - MARGIN - rxx
    # banco
    card(s, MARGIN, 2.7, lw, 2.75, fill=NAVY, line=NAVY, radius=0.05, shadow=True)
    txt(s, MARGIN + 0.32, 2.95, lw - 0.6, 0.3, "O BANCO", 10, font=MONO, color=GOLD, char=1.6, bold=True)
    txt(s, MARGIN + 0.32, 3.28, lw - 0.6, 0.4, "Administrador fiduciário formal", 15, font=SERIF, color=D_TXT, bold=True)
    for i, t in enumerate(["A licença e a estrutura que já existem", "Os diretores estatutários responsáveis",
                           "A responsabilidade perante a CVM", "A conta e o nome"]):
        yy = 3.78 + i * 0.42
        oval(s, MARGIN + 0.4, yy + 0.1, 0.035, fill=GOLD)
        txt(s, MARGIN + 0.58, yy, lw - 0.9, 0.4, t, 11, font=SANS, color=D_TXT, spacing=1.05)
    # startup
    card(s, rxx, 2.7, rw, 2.75, fill=CARD, line=LINE, radius=0.05, shadow=True)
    txt(s, rxx + 0.32, 2.95, rw - 0.6, 0.3, "A ARGUS (NOSSA EQUIPE)", 10, font=MONO, color=NAVY, char=1.4, bold=True)
    txt(s, rxx + 0.32, 3.28, rw - 0.6, 0.4, "Tecnologia + operação", 15, font=SERIF, color=INK, bold=True)
    for i, t in enumerate(["A plataforma completa, já construída", "A operação da controladoria, ponta a ponta",
                           "A originação de novos fundos", "Os manuais e controles que protegem o banco"]):
        yy = 3.78 + i * 0.42
        oval(s, rxx + 0.4, yy + 0.1, 0.035, fill=NAVY)
        txt(s, rxx + 0.58, yy, rw - 0.9, 0.4, t, 11, font=SANS, color=INK, spacing=1.05)
    # split
    rect(s, MARGIN, 5.7, PW - 2 * MARGIN, 0.95, fill=SOFT, line=LINE, radius=0.05, shadow=False)
    txt(s, MARGIN + 0.35, 5.83, 5.9, 0.35, "Divisão da receita de administração", 13, font=SERIF, color=INK, bold=True)
    txt(s, MARGIN + 0.35, 6.17, 6.05, 0.5, "O banco (administrador) remunera a Argus pelos serviços de tecnologia "
        "e operação. Começa 50/50 — com opção de recompra que pode levar a 15/85.", 9.5, font=SANS, color=MUTED, spacing=1.12)
    txt(s, PW - MARGIN - 3.5, 5.8, 3.2, 0.6, "50 / 50", 30, font=SERIF, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
    txt(s, PW - MARGIN - 3.5, 6.4, 3.2, 0.3, "BANCO / ARGUS — INÍCIO", 8, font=MONO, color=MUTED, char=1.0, align=PP_ALIGN.RIGHT)
    source(s, "Res. CVM 175 art. 118 §1º (repasse a prestador contratado). Linha vermelha: o banco é o administrador de verdade — decide, controla e responde.")

def s_conta_banco():
    s = slide(PAPER)
    head(s, "A parceria", "A conta para o banco", "Receita nova, custo de entrada quase zero.")
    txt(s, 8.5, 2.6, PW - MARGIN - 8.5, 0.3, "FATIA INICIAL DO BANCO (50%) — R$ MIL/ANO", 8.5, font=MONO, color=MUTED, char=0.8)
    bar_chart(s, 8.35, 2.95, PW - MARGIN - 8.35, 3.35,
              ["R$50mi", "R$200mi", "R$500mi", "R$1bi"],
              [("R$ mil", [20, 80, 200, 350])],
              colors=[GOLD], col=True, show_val=True, val_fmt='0', grid=False, gap=55, val_color=NAVY)
    txt(s, 8.35, 6.32, 4.3, 0.4, "PL sob administração. 50% de 0,08% a.a.\n(participação inicial; ver opção de recompra).",
        8.5, font=MONO, color=FAINT, char=0.2, spacing=1.15)
    items = [
        ("Custo de entrada ~zero", "sem investir em tecnologia nem contratar equipe — a receita nasce sobre a estrutura atual."),
        ("Dispensa de capital mínimo", "por já ser instituição autorizada, o banco não precisa dos R$ 550 mil de capital exigidos de uma PJ comum."),
        ("Alavancagem operacional", "a estrutura custa ~R$ 27 mil/ano fixos; quase toda receita nova vira resultado."),
        ("Originação incluída", "somos nós que trazemos os fundos — o banco não precisa de esforço comercial."),
    ]
    y = 2.62
    for i, (t, d) in enumerate(items):
        oval(s, MARGIN + 0.06, y + 0.12, 0.05, fill=GOLD)
        txt(s, MARGIN + 0.26, y, 7.2, 0.32, t, 13, font=SERIF, color=INK, bold=True)
        txt(s, MARGIN + 0.26, y + 0.34, 7.2, 0.55, d, 10.5, font=SANS, color=MUTED, spacing=1.18)
        y += 0.82
    rect(s, MARGIN, 5.88, 7.25, 0.66, fill=SOFT, line=LINE, radius=0.06, shadow=False)
    txt(s, MARGIN + 0.22, 5.99, 6.85, 0.5, "A explorar com a tesouraria: o caixa dos fundos em conta "
        "(float) e a compra de CDB do banco pelos fundos de RF (funding) — upside a quantificar.",
        9.5, font=SANS, color=INK, italic=True, spacing=1.12)
    source(s, "guia_potencial_financeiro.md; guia_credenciamento_banco.md (dispensa de capital). Receita = PL × taxa; não é promessa de PL.")

def s_vigilancia_porque():
    s = slide(NAVY_DK)
    head(s, "A vigilância", "Por que vigiar cada fundo", "O maior medo do banco tem nome e valor.", dark=True)
    txt(s, MARGIN, 2.5, PW - 2 * MARGIN, 0.55,
        "A CVM pune o administrador e o custodiante quando falham na diligência — não só o fraudador. "
        "Cada fundo é uma porta que alguém pode tentar usar para fraude ou lavagem.",
        13, font=SANS, color=D_MUT, spacing=1.3)
    casos = [
        ("Silverado · FIDCs", "2024",
         "Recebíveis fictícios; cedentes ligados à gestora. Os administradores foram multados por falha de diligência.",
         "R$ 490 mi · BNY Mellon, Santander e Gradual punidos"),
        ("Banco Master · fundos", "2025",
         "Fundos usados para triangular caixa e inflar ativos — R$ 850 mi marcados como R$ 10 bi. Banco liquidado.",
         "veículo de lavagem · maior acionamento do FGC"),
        ("FIP LSH", "2024",
         "Laudo do único ativo superfaturado; transferência de riqueza via cotas. Precificação de ilíquido como vetor.",
         "gestor multado em R$ 20 mi"),
    ]
    cw = (PW - 2 * MARGIN - 2 * 0.4) / 3
    for i, (t, ano, d, ev) in enumerate(casos):
        x = MARGIN + i * (cw + 0.4)
        card(s, x, 3.35, cw, 2.35, fill=NAVY_DP, line=D_LINE, radius=0.05, shadow=False)
        txt(s, x + 0.28, 3.54, cw - 1.05, 0.3, t, 12.5, font=SERIF, color=D_TXT, bold=True)
        txt(s, x + cw - 0.82, 3.57, 0.6, 0.28, ano, 9, font=MONO, color=NEG, char=0.4, align=PP_ALIGN.RIGHT)
        txt(s, x + 0.28, 3.94, cw - 0.56, 1.05, d, 10, font=SANS, color=D_MUT, spacing=1.22)
        hline(s, x + 0.28, 5.02, cw - 0.56, color=D_LINE, weight=0.75)
        txt(s, x + 0.28, 5.12, cw - 0.56, 0.5, ev, 9, font=MONO, color=GOLD, char=0.2, spacing=1.18)
    rect(s, MARGIN, 5.95, PW - 2 * MARGIN, 0.66, fill=NAVY_DP, line=D_LINE, radius=0.05, shadow=False)
    txt(s, MARGIN + 0.28, 6.08, PW - 2 * MARGIN - 0.56, 0.45,
        "Os vetores se repetem: marcação inflada · ativo sem lastro · partes relacionadas · lavagem via cotas · descasamento de liquidez.",
        11, font=SANS, color=D_TXT, italic=True, spacing=1.1)
    source(s, "CVM — PAS Silverado 19957.006858/2019-25 (2024); imprensa (Banco Master, 2025); CVM (FIP LSH).", dark=True)

def s_vigilancia_como():
    s = slide(PAPER)
    head(s, "A vigilância", "Como vigiamos barato", "O mesmo motor roda para 1 ou 500 fundos.")
    controles = [
        ("Conciliação diária", "ativo fantasma — posição não bate com o custodiante"),
        ("Preço fora da curva", "marcação inflada — desvio da referência ANBIMA/B3"),
        ("Grafo de partes relacionadas", "autonegociação — contraparte ligada ao gestor"),
        ("Anomalia e timing", "cotização e lavagem — aplica-resgata, padrão atípico"),
        ("Enquadramento contínuo", "desenquadramento — limite violado no dia"),
        ("KYC · PLD · COAF", "beneficiário oculto — 24 h para comunicar"),
    ]
    lx = MARGIN; lw = 6.4
    txt(s, lx, 2.55, lw, 0.28, "O QUE RODA TODO DIA, SOZINHO", 9, font=MONO, color=NAVY, char=1.4, bold=True)
    for i, (c, v) in enumerate(controles):
        yy = 2.94 + i * 0.49
        oval(s, lx + 0.06, yy + 0.12, 0.045, fill=GOLD)
        txt(s, lx + 0.24, yy, 2.75, 0.35, c, 11.5, font=SERIF, color=INK, bold=True)
        txt(s, lx + 3.0, yy + 0.03, lw - 3.0, 0.35, v, 9.5, font=SANS, color=MUTED, spacing=1.1)
    rx = 7.7; rw = PW - MARGIN - rx
    card(s, rx, 2.9, rw, 2.5, fill=NAVY, line=NAVY, radius=0.05)
    txt(s, rx + 0.3, 3.1, rw - 0.6, 0.3, "O ARGUMENTO PARA O BANCO", 9, font=MONO, color=GOLD, char=1.2, bold=True)
    txt(s, rx + 0.3, 3.46, rw - 0.6, 1.0,
        "As grandes não vigiam fundos pequenos assim — o custo manual não paga. Nós vigiamos todos com o "
        "mesmo rigor, porque o custo é de software.", 11.5, font=SANS, color=D_TXT, spacing=1.24)
    hline(s, rx + 0.3, 4.62, rw - 0.6, color=D_LINE)
    txt(s, rx + 0.3, 4.76, rw - 0.6, 0.55,
        "O fundo pequeno fica MAIS bem vigiado que na administradora tradicional — é o que nos deixa baratos "
        "e seguros ao mesmo tempo.", 11, font=SERIF, color=GOLD, bold=True, spacing=1.16)
    txt(s, MARGIN, 5.85, PW - 2 * MARGIN, 0.55,
        "Honesto: a IA prioriza alertas e grava a trilha de auditoria — não decide sozinha. Lastro real, conluio e "
        "documento forjado ainda exigem verificação humana/forense. No protótipo: 7 regras (R1–R7), grafo de vínculos e triagem com trilha.",
        9.5, font=SANS, color=MUTED, italic=True, spacing=1.18)
    source(s, "guia_controle_fraudes_riscos.md; Res. CVM 50/2021 (PLD, arts. 20 e 22); admin/fraude.php do protótipo.")

def s_risco_banco():
    s = slide(NAVY_DK)
    head(s, "A parceria", "O risco do banco", "“E quando um gestor fizer besteira?” — três camadas de blindagem.", dark=True)
    layers = [
        ("A norma protege", "Res. CVM 175: responsabilidade subjetiva (exige dolo/culpa) e de meio, não de fim. Sem solidariedade automática — cada um responde pela sua esfera."),
        ("A operação produz prova", "nenhuma cota publica sem aprovação formal do gestor; toda correção vira lançamento com trilha; conciliação diária em três frentes."),
        ("O monitoramento é total", "IA lê regulamentos e vigia preços fora da curva, partes relacionadas e ativos sem lastro — em 100% dos fundos, todos os dias."),
        ("O compliance manda", "o compliance do banco designa o diretor, aprova os manuais, recebe relatórios e audita quando quiser. Envolvimento é blindagem."),
    ]
    x0 = MARGIN; cw = (PW - 2 * MARGIN - 3 * 0.35) / 4
    for i, (t, d) in enumerate(layers):
        x = x0 + i * (cw + 0.35)
        card(s, x, 2.75, cw, 3.55, fill=NAVY_DP, line=D_LINE, radius=0.05, shadow=False)
        txt(s, x + 0.24, 3.0, cw - 0.48, 0.35, f"{i+1:02d}", 13, font=MONO, color=GOLD, bold=True)
        txt(s, x + 0.24, 3.42, cw - 0.48, 0.7, t, 14, font=SERIF, color=D_TXT, bold=True, spacing=1.02)
        hline(s, x + 0.24, 4.12, cw - 0.48, color=D_LINE)
        txt(s, x + 0.24, 4.28, cw - 0.48, 1.9, d, 10, font=SANS, color=D_MUT, spacing=1.26)
    source(s, "guia_obrigacoes_riscos_banco.md; Res. CVM 175 art. 81.", dark=True)

def s_piloto():
    s = slide(PAPER)
    head(s, "A prova", "Não é ideia em guardanapo", "Um protótipo operacional, de ponta a ponta.")
    txt(s, MARGIN, 2.55, 5.1, 0.55, "argusdtvm.com.br", 19, font=MONO, color=NAVY, bold=True, char=0.5)
    items = [
        "Do checklist de abertura ao informe diário à CVM.",
        "Batch diário: marcação, cota e PL de qualquer data.",
        "Aprovação de cota pelo gestor e trilha de correções.",
        "Custódia, conciliação e contabilidade em partidas dobradas.",
        "Passivo do cotista: come-cotas, IR e IOF (Lei 14.754).",
    ]
    y = 3.3
    for it in items:
        oval(s, MARGIN + 0.06, y + 0.11, 0.04, fill=GOLD)
        txt(s, MARGIN + 0.24, y, 4.95, 0.55, it, 11, font=SANS, color=INK, spacing=1.2)
        y += 0.58
    rect(s, MARGIN, 6.2, 5.1, 0.66, fill=SOFT, line=LINE, radius=0.06, shadow=False)
    txt(s, MARGIN + 0.22, 6.31, 4.7, 0.5, "Nota honesta: dados simulados. O protótipo prova a capacidade "
        "operacional — não substitui a licença nem a parceria que propomos.", 8.5, font=SANS, color=MUTED, italic=True, spacing=1.12)
    mx, my, mw, mh = 6.7, 2.55, PW - MARGIN - 6.7, 4.05
    card(s, mx, my, mw, mh, fill=NAVY_DK, line=NAVY_DK, radius=0.03, shadow=True)
    for i in range(3):
        oval(s, mx + 0.32 + i * 0.18, my + 0.3, 0.045, fill=RGBColor(0x3a, 0x4a, 0x5e))
    txt(s, mx + 1.1, my + 0.16, mw - 1.3, 0.3, "ARGUS · PAINEL DO ADMINISTRADOR", 8, font=MONO, color=D_MUT, char=1.0)
    hline(s, mx + 0.3, my + 0.52, mw - 0.6, color=D_LINE, weight=0.5)
    kpis = [("FUNDOS ATIVOS", "9"), ("COTA D-0", "1,042318"), ("COTISTAS", "130"), ("FECHAMENTO", "9/9 OK")]
    kw = (mw - 0.6 - 0.3) / 2
    for i, (l, v) in enumerate(kpis):
        col = i % 2; row = i // 2
        kx = mx + 0.3 + col * (kw + 0.3); ky = my + 0.72 + row * 0.9
        rect(s, kx, ky, kw, 0.78, fill=NAVY_DP, line=D_LINE, radius=0.06, shadow=False)
        txt(s, kx + 0.2, ky + 0.13, kw - 0.4, 0.3, l, 7.5, font=MONO, color=D_MUT, char=1.0)
        txt(s, kx + 0.2, ky + 0.34, kw - 0.4, 0.4, v, 16, font=SERIF, color=D_TXT, bold=True)
    txt(s, mx + 0.3, my + mh - 0.5, mw - 0.6, 0.3, "Informe diário à CVM · gerado — D+1", 8.5, font=MONO, color=GOLD, char=0.5)

def s_regulatorio():
    s = slide(PAPER)
    head(s, "O caminho", "Mapa regulatório", "Todo o caminho já levantado e documentado.")
    txt(s, MARGIN, 2.55, PW - 2 * MARGIN, 0.5,
        "Nós preparamos os dossiês; o banco protocola e supervisiona. Nada aqui é território desconhecido.",
        13, font=SANS, color=INK, spacing=1.25)
    docs = [
        ("Res. CVM 21", "Administrador fiduciário", "Habilitação do banco (SSM → ANBIMA → CVM). Prazo legal de 60 dias. Formulário, manuais e website entregues prontos."),
        ("Res. CVM 32", "Custodiante", "Autorização separada. Pode ser terceirizada no início e internalizada em escala. Exige auditoria interna — não auditor externo."),
        ("Adesões de mercado", "B3 · SELIC · RSFN", "Contas segregadas, liquidação DVP e mensageria do SPB. Ocorrem após o deferimento da custódia."),
        ("Dossiês prontos", "Pacotes completos", "Matrizes de conformidade e manuais em custodiante/ e administrador_fiduciario/ — com o protótipo como anexo técnico."),
    ]
    x0 = MARGIN; cw = (PW - 2 * MARGIN - 3 * 0.4) / 4
    for i, (r, t, d) in enumerate(docs):
        x = x0 + i * (cw + 0.4)
        card(s, x, 3.2, cw, 2.95, radius=0.04)
        txt(s, x + 0.26, 3.44, cw - 0.5, 0.32, r, 12, font=MONO, color=GOLD, bold=True, char=0.5)
        txt(s, x + 0.26, 3.8, cw - 0.5, 0.6, t, 13.5, font=SERIF, color=INK, bold=True, spacing=1.02)
        hline(s, x + 0.26, 4.42, cw - 0.52, color=LINE)
        txt(s, x + 0.26, 4.56, cw - 0.5, 1.55, d, 10, font=SANS, color=MUTED, spacing=1.24)
    source(s, "guia_credenciamento_banco.md; custodiante/README.md; administrador_fiduciario/README.md.")

def s_roadmap():
    s = slide(PAPER)
    head(s, "O caminho", "Cronograma", "Do acordo ao go-live, em 12 a 18 meses.")
    phases = [
        ("Acordo", "Carta de intenções e consulta jurídica conjunta.", "Mês 0"),
        ("Habilitação", "Administrador fiduciário na CVM (Res. 21, ~60 dias).", "Meses 1–4"),
        ("Custódia", "Res. 32 + adesões B3/SELIC/RSFN, em paralelo.", "Meses 3–9"),
        ("Homologação", "Mensageria, manuais e controles validados.", "Meses 6–12"),
        ("Operação assistida", "3 a 5 fundos-piloto sob a licença, antes de escalar.", "Meses 12–18"),
    ]
    n = len(phases); y = 3.5
    x0 = MARGIN + 1.0; span = PW - 2 * MARGIN - 2.0; step = span / (n - 1); Lw = 2.25
    hline(s, x0, y, span, color=LINE, weight=1.2)
    for i, (t, d, when) in enumerate(phases):
        cx = x0 + i * step
        oval(s, cx, y, 0.09, fill=NAVY, line=PAPER, line_w=2.0); oval(s, cx, y, 0.035, fill=GOLD)
        lxb = cx - Lw / 2
        txt(s, lxb, y - 0.62, Lw, 0.3, when, 9, font=MONO, color=GOLD, char=0.6, align=PP_ALIGN.CENTER)
        txt(s, lxb, y + 0.28, Lw, 0.4, t, 13, font=SERIF, color=INK, bold=True, align=PP_ALIGN.CENTER)
        txt(s, lxb, y + 0.66, Lw, 1.1, d, 9.5, font=SANS, color=MUTED, spacing=1.18, align=PP_ALIGN.CENTER)
    rect(s, MARGIN, 6.05, PW - 2 * MARGIN, 0.62, fill=SOFT, line=LINE, radius=0.06, shadow=False)
    txt(s, MARGIN + 0.25, 6.17, PW - 2 * MARGIN - 0.5, 0.45,
        "Começamos pela administração (viável já com poucos fundos) e tratamos a custódia em paralelo, "
        "internalizando-a só quando o volume de fundos justificar.", 10, font=SANS, color=MUTED, italic=True, spacing=1.12)
    source(s, "guia_custodia_conexoes.md; relatorio_gaps_producao_licenca.md. Prazos indicativos; o legal da Res. 21 é 60 dias.")

def s_riscos():
    s = slide(PAPER)
    head(s, "Honestidade", "Riscos", "Os riscos, antes que vocês perguntem.")
    table(s, MARGIN, 2.75, [4.0, 3.6, PW - 2 * MARGIN - 4.0 - 3.6],
          ["RISCO", "POR QUE IMPORTA", "COMO ENFRENTAMOS"],
          [["Custo por fundo irredutível", "parte do custo pode não ser automatizável", "o piloto mede o custo unitário real antes de escalar"],
           ["IA não pega fraudador competente", "a fraude é feita para as contas fecharem", "IA é apoio de uma camada humana de compliance, no plano"],
           ["Auditor único barato", "compromete a independência exigida", "pulverizamos entre vários auditores"],
           ["Receita modesta no início", "a taxa é enxuta de propósito", "custo de entrada do banco ~zero e alavancagem operacional"]],
          rh=0.66, accent_col=2, csize=10.5)
    txt(s, MARGIN, 6.35, PW - 2 * MARGIN, 0.5,
        "Se algum desses riscos matar a tese, queremos descobrir no piloto — não na escala. É para isso "
        "que a fase assistida existe.", 11, font=SANS, color=MUTED, italic=True, spacing=1.15)
    source(s, "plano_dtvm_fundos_pequenos.md (§8).")

def s_pedido():
    s = slide(NAVY_DK)
    head(s, "O pedido", "Próximos passos", "O que propomos hoje — três passos concretos.", dark=True)
    steps = [
        ("01", "Reunião técnica", "com o time de tecnologia e o compliance de vocês, para abrir o modelo, os controles e o protótipo."),
        ("02", "Carta de intenções", "não-vinculante, para detalharmos a divisão de resultado e de responsabilidades."),
        ("03", "3 a 5 fundos-piloto", "sob a licença do banco, com supervisão total de vocês e o custo unitário medido."),
    ]
    x0 = MARGIN; cw = (PW - 2 * MARGIN - 2 * 0.45) / 3
    for i, (n, t, d) in enumerate(steps):
        x = x0 + i * (cw + 0.45)
        card(s, x, 2.9, cw, 3.0, fill=NAVY_DP, line=D_LINE, radius=0.05, shadow=False)
        txt(s, x + 0.3, 3.18, cw - 0.6, 0.5, n, 20, font=MONO, color=GOLD, bold=True)
        txt(s, x + 0.3, 3.78, cw - 0.6, 0.5, t, 17, font=SERIF, color=D_TXT, bold=True)
        hline(s, x + 0.3, 4.32, cw - 0.6, color=D_LINE)
        txt(s, x + 0.3, 4.48, cw - 0.6, 1.3, d, 11, font=SANS, color=D_MUT, spacing=1.28)
    txt(s, MARGIN, 6.25, PW - 2 * MARGIN, 0.5,
        "Vocês já têm o mais difícil — a licença e a estrutura. Nós trazemos tudo pronto e reduzimos o "
        "seu risco. Podemos marcar a reunião técnica?", 12.5, font=SERIF, color=GOLD, italic=True, spacing=1.2)

def s_fecho():
    s = slide(NAVY_DK)
    draw_eye(s, PW / 2, 2.5, 1.4, outline=D_TXT, ring=TEAL, pupil=GOLD, highlight=True)
    txt(s, 0, 3.55, PW, 0.7, "ARGUS", 40, font=SERIF, color=D_TXT, bold=True, char=1.5, align=PP_ALIGN.CENTER)
    txt(s, 0, 4.35, PW, 0.5, "A administradora desenhada para os fundos que ninguém quer atender.", 15,
        font=SERIF, color=GOLD, italic=True, align=PP_ALIGN.CENTER)
    hline(s, PW / 2 - 1.6, 5.15, 3.2, color=D_LINE)
    txt(s, 0, 5.4, PW, 0.4, "argusdtvm.com.br", 13, font=MONO, color=D_MUT, char=1.5, align=PP_ALIGN.CENTER)
    txt(s, 0, PH - 0.9, PW, 0.3, "PROPOSTA DE PARCERIA · CONFIDENCIAL · PROTÓTIPO COM DADOS SIMULADOS", 8.5,
        font=MONO, color=D_FAINT, char=1.6, align=PP_ALIGN.CENTER)

# =================================================================== ANEXO
def s_anexo_divider():
    # Sumário do material de apoio — vem DEPOIS do slide final (s_fecho) da apresentação principal.
    s = slide(NAVY_DK)
    draw_eye(s, MARGIN + 0.13, 0.62, 0.30)
    txt(s, MARGIN + 0.34, 0.47, 3, 0.32, "ARGUS", 12.5, font=SERIF, color=D_TXT, bold=True, char=0.5)
    hline(s, MARGIN, 0.92, PW - 2 * MARGIN, color=D_LINE)
    kicker(s, MARGIN, 1.2, "Anexo · sumário")
    txt(s, MARGIN, 1.54, 11, 0.9, "Material de apoio", 30, font=SERIF, color=D_TXT)
    txt(s, MARGIN, 2.42, 10.6, 0.4, "A apresentação principal termina no slide anterior. Os slides a seguir respondem às perguntas prováveis.",
        11.5, font=SANS, color=D_MUT, spacing=1.2)
    grupos_l = [
        ("PRODUTO", ["Viabilidade do cotista (% do CDI)", "Custo de cada fundo"]),
        ("NEGÓCIO", ["Receita por escala", "Custo da estrutura"]),
        ("MERCADO", ["Concorrência", "O tamanho da cauda", "O piso dos grandes"]),
    ]
    grupos_r = [
        ("VIGILÂNCIA", ["Tipologia → sinal → limite", "Casos reais de fraude"]),
        ("REGULATÓRIO", ["Taxa CVM do fundo", "Quem faz o quê", "Obrigações de PLD/FT"]),
        ("PARCERIA", ["O que o banco faz"]),
        ("MÉTODO", ["Fontes"]),
    ]
    for ci, (cx, gs) in enumerate([(MARGIN, grupos_l), (6.95, grupos_r)]):
        yy = 2.98
        for cat, items in gs:
            txt(s, cx, yy, 5.4, 0.28, cat, 9, font=MONO, color=GOLD, char=1.6, bold=True)
            yy += 0.30
            for it in items:
                oval(s, cx + 0.06, yy + 0.11, 0.035, fill=TEAL)
                txt(s, cx + 0.22, yy, 5.2, 0.28, it, 11.5, font=SANS, color=D_TXT)
                yy += 0.27
            yy += 0.12
    hline(s, MARGIN, PH - 0.86, PW - 2 * MARGIN, color=D_LINE)
    txt(s, MARGIN, PH - 0.72, 9, 0.3, "ANEXO · DADOS, VIGILÂNCIA E FONTES", 8.5, font=MONO, color=D_FAINT, char=1.8)

def s_ax_cotista():
    s = slide(PAPER)
    head(s, "Anexo · produto", "Viabilidade do cotista", "% do CDI — 1º ano e regime permanente.")
    line_chart(s, MARGIN, 2.9, 7.3, 3.5,
               ["R$1mi", "R$2mi", "R$5mi", "R$10mi", "R$50mi"],
               [("1º ano (sem taxa CVM)", [100.3, 101.4, 101.8, 102.0, 102.1]),
                ("Regime permanente", [97.3, 99.9, 101.2, 101.5, 101.9]),
                ("CDI", [100, 100, 100, 100, 100])],
               colors=[TEAL, NAVY, FAINT], smooth=False, val_fmt='0.0"%"', legend=True, widths=[2.2, 2.6, 1.0])
    txt(s, 8.6, 3.0, PW - MARGIN - 8.6, 3.2,
        "No 1º ano o fundo não paga taxa CVM (se criado após abril) e fecha bem já em R$ 1 milhão. "
        "No regime permanente, a taxa CVM entra e o fundo de R$ 1 mi vira caso-limite (97% do CDI). "
        "A partir de R$ 5 mi, supera o CDI nos dois regimes.", 12, font=SANS, color=INK, spacing=1.34)
    source(s, "planilha_custos.md (§7.1 e §7.2). Não vender o número do 1º ano como permanente.")

def s_ax_custo_fundo():
    s = slide(PAPER)
    head(s, "Anexo · produto", "Custo de cada fundo", "Onde vai o custo de um fundo de R$ 5 milhões.")
    stacked_col(s, MARGIN, 2.95, 5.6, 3.4, ["Custo anual\ndo fundo"],
                [("Gestão · 0,70% (do gestor)", [0.70]), ("Taxa CVM · R$3.162", [0.06]),
                 ("Auditoria · R$1.500", [0.03]), ("Administração · 0,08% (Argus)", [0.08]),
                 ("Custódia · R$0 (banco)", [0.001])],
                colors=[MUTED, NEG, GOLD_DK, NAVY, TEAL], val_fmt='0.00"%"', show_val=False)
    txt(s, 6.9, 2.95, PW - MARGIN - 6.9, 3.4,
        "A fatia da Argus é 0,08% — a menor da conta. O peso está na gestão (do gestor) e nos custos "
        "fixos regulatórios (taxa CVM, auditoria), que a plataforma reduz ao mínimo e a custódia via "
        "banco zera para o fundo. É por isso que o fundo pequeno passa a fechar.", 12, font=SANS, color=INK, spacing=1.34)
    source(s, "planilha_custos.md (§7). Taxa CVM R$3.162/ano (faixa até R$5mi); auditoria adaptada em lote.")

def s_ax_receita_pl():
    s = slide(PAPER)
    head(s, "Anexo · negócio", "Receita por escala", "A receita = PL × taxa. Cresce com a base.")
    txt(s, MARGIN, 2.6, PW - 2 * MARGIN, 0.3, "RECEITA BRUTA ANUAL DA ADMINISTRADORA — R$ MIL", 9, font=MONO, color=MUTED, char=1.0)
    bar_chart(s, MARGIN, 2.95, PW - 2 * MARGIN, 3.4,
              ["R$5mi", "R$50mi", "R$200mi", "R$500mi", "R$800mi"],
              [("Taxa 0,08%", [4, 40, 160, 400, 640]), ("Taxa 0,05% (grandes)", [2.5, 25, 100, 250, 400])],
              colors=[NAVY, TEAL], col=True, show_val=True, val_fmt='0', grid=True, gap=45, legend=True, val_color=MUTED)
    source(s, "guia_potencial_financeiro.md (§3). Fundos até ~R$800mi são o sweet spot: custo marginal ~zero permite taxa agressiva. Arit. sobre o PL — não é projeção de PL.")

def s_ax_estrutura():
    s = slide(PAPER)
    head(s, "Anexo · negócio", "Custo da estrutura", "Operar a administradora custa pouco.")
    table(s, MARGIN, 2.75, [5.0, 3.2, PW - 2 * MARGIN - 5.0 - 3.2],
          ["ITEM", "VALOR/ANO", "NATUREZA"],
          [["Nuvem (AWS)", "~R$ 5.400", "escalável; ~R$0 sob o banco"],
           ["Taxa CVM do administrador (PJ)", "~R$ 9.519", "some se operar sob o banco"],
           ["Contribuição ANBIMA", "negociada", "some se o banco já é associado"],
           ["Seguro E&O, contador, jurídico", "a cotar", "recorrente / pontual"]],
          rh=0.54, accent_col=1, csize=10.5)
    rect(s, MARGIN, 5.78, PW - 2 * MARGIN, 0.82, fill=SOFT, line=LINE, radius=0.05, shadow=False)
    txt(s, MARGIN + 0.3, 5.9, PW - 2 * MARGIN - 0.6, 0.6,
        "Total da estrutura: da ordem de R$ 27 mil/ano (custódia terceirizada). Como o custo fixo é baixo "
        "e a receita escala com o PL, dobrar a base quase dobra o resultado — alavancagem operacional "
        "típica de software. O gargalo nunca foi custo; é fechar a parceria e captar fundos.",
        10.5, font=SANS, color=INK, spacing=1.22)
    source(s, "planilha_custos.md (§2–6); guia_potencial_financeiro.md (§2).")

def s_ax_taxacvm():
    s = slide(PAPER)
    head(s, "Anexo · regulatório", "Taxa CVM do fundo", "O “vilão” dos fundos minúsculos — fixo por faixa.")
    bar_chart(s, MARGIN, 2.95, 7.2, 3.4,
              ["até R$2,5mi", "R$2,5–5mi", "R$5–10mi", "R$10–20mi", "R$20–40mi"],
              [("R$/ano", [2400, 3600, 5400, 7200, 9600])],
              colors=[NAVY], col=True, show_val=True, val_fmt='#,##0', grid=False, gap=50, val_color=NAVY)
    txt(s, 8.5, 3.0, PW - MARGIN - 8.5, 3.3,
        "A taxa de fiscalização da CVM é fixa por faixa de PL — não cai para o fundo pequeno. É paga pelo "
        "fundo (não pela administradora) e some no 1º ano de fundos criados após abril.\n\n"
        "Não é diluída por subclasse: incide sobre o PL, fatiado ou não.", 12, font=SANS, color=INK, spacing=1.32)
    source(s, "Tabela oficial CVM (Lei 14.317). Valores anuais aproximados; confirmar vigência 2026 com a CVM.")

def s_ax_concorrencia():
    s = slide(PAPER)
    head(s, "Anexo · mercado", "Concorrência", "A infraestrutura de fundos foi toda para o FIDC.")
    table(s, MARGIN, 2.8, [3.0, 5.4, PW - 2 * MARGIN - 3.0 - 5.4],
          ["PLAYER", "FOCO CONFIRMADO", "FUNDO PEQUENO GENÉRICO?"],
          [["Kanastra", "FIDC / crédito estruturado (exclusivo)", "Não — explicitamente fora"],
           ["QI Tech / Singulare", "Infra de crédito (SCD) + FIDC (>1.000 fundos)", "Não"],
           ["Vórtx", "FIDC + dívida estruturada (Full FIDC)", "Só gestoras grandes"],
           ["BRLTrust", "FIDC + crédito corporativo/distressed", "Não como foco"],
           ["Singulare", "FIDC (líder Uqbar há 10+ anos)", "Não"]],
          rh=0.6, accent_col=2, csize=10)
    txt(s, MARGIN, 6.2, PW - 2 * MARGIN, 0.5,
        "Todas convergiram para crédito estruturado de ticket alto. O nicho de fundos pequenos genéricos "
        "(multimercado e RF) segue sem um dono — é exatamente onde entramos.", 11, font=SANS, color=MUTED, italic=True, spacing=1.15)
    source(s, "Sites das empresas e imprensa (Finsiders, Bloomberg Línea, NeoFeed, Let's Money), 2024–2026.")

def s_ax_mercado():
    s = slide(PAPER)
    head(s, "Anexo · mercado", "O tamanho da cauda", "Muitos fundos pequenos, muitas gestoras pequenas.")
    stats = [("~31 mil", "fundos de investimento no país (dez/2024)"),
             ("1.277", "fundos com PL acima de R$ 1 bilhão — o resto é a cauda"),
             ("976", "gestoras; 51% com menos de R$ 400 milhões sob gestão"),
             ("27.515", "assessores de investimento (+502% em uma década)")]
    x0 = MARGIN; cw = (PW - 2 * MARGIN - 3 * 0.4) / 4
    for i, (v, l) in enumerate(stats):
        x = x0 + i * (cw + 0.4)
        txt(s, x, 2.9, cw, 0.6, v, 27, font=SERIF, color=NAVY, bold=True)
        hline(s, x, 3.56, cw - 0.15, color=GOLD, weight=1.6)
        txt(s, x, 3.7, cw, 1.4, l, 11, font=SANS, color=MUTED, spacing=1.26)
    hline(s, MARGIN, 5.4, PW - 2 * MARGIN, color=LINE)
    txt(s, MARGIN, 5.6, PW - 2 * MARGIN, 1.0,
        "A concentração é no topo: as 50 maiores gestoras detêm 83% do patrimônio. A base — milhares de "
        "fundos pequenos e centenas de casas pequenas — é numerosa e mal atendida. É um mercado de muitos "
        "clientes de ticket baixo, que só fecha com custo de software.", 12.5, font=SANS, color=INK, spacing=1.3)
    source(s, "ANBIMA (2024); ANCORD (assessores, dez/2025); imprensa especializada.")

def s_ax_precos():
    s = slide(PAPER)
    head(s, "Anexo · mercado", "O piso dos grandes", "Taxa de administração em regulamentos reais.")
    table(s, MARGIN, 2.85, [3.5, 4.0, PW - 2 * MARGIN - 3.5 - 4.0],
          ["ADMINISTRADOR", "TAXA DE ADMINISTRAÇÃO", "PISO MENSAL"],
          [["BRL Trust", "0,30% a.a. (ou o piso, o maior)", "R$ 25.000"],
           ["Vórtx", "1,25% a.a.", "R$ 25.000"],
           ["Reag", "0,30% a.a.", "R$ 20.000"],
           ["Oliveira Trust", "0,06–0,16% a.a. (regressiva)", "R$ 16.250"],
           ["Singulare (hoje QI Tech)", "pacote adm + custódia", "R$ 15.500"],
           ["Intrag / Itaú", "1,20% a.a. + piso", "R$ 3.000"],
           ["Argus", "0,08% a.a. + piso", "R$ 100"]],
          rh=0.4, accent_col=2, csize=10, dsize=11.5)
    source(s, "Regulamentos reais na CVM/Fundos.NET (fundos específicos, majoritariamente estruturados). Fintechs (Kanastra) 0,1–0,5% s/ ativos. Ver benchmark_precos_concorrencia.md.")

def s_ax_papeis():
    s = slide(PAPER)
    head(s, "Anexo · regulatório", "Quem faz o quê", "Os papéis de um fundo — e onde entramos.")
    table(s, MARGIN, 2.8, [3.6, 5.3, PW - 2 * MARGIN - 3.6 - 5.3],
          ["FUNÇÃO", "RESPONSABILIDADE", "NA PARCERIA"],
          [["Administrador fiduciário", "constituição, cota, PL, contabilidade e informes", "Banco responde · Argus executa"],
           ["Gestor", "decisões de investimento e enquadramento", "cliente ou parceiro"],
           ["Custodiante", "guarda dos ativos, liquidação e conciliação", "banco (ou terceiro no início)"],
           ["Distribuidor", "coloca o investidor dentro do fundo", "banco — sem licença nova"]],
          rh=0.68, accent_col=2, csize=10.5)
    txt(s, MARGIN, 6.25, PW - 2 * MARGIN, 0.5,
        "Distribuir cotas dos fundos que ele mesmo administra quase sempre já está no escopo do banco "
        "(Lei 6.385, art. 15) — um eixo extra de valor, sem obter autorização nova.", 11, font=SANS, color=MUTED, italic=True, spacing=1.15)
    source(s, "Res. CVM 21/32/175; distribuidor/manual_banco_distribuidor.md.")

def s_ax_monitoramento():
    s = slide(PAPER)
    head(s, "Anexo · vigilância", "O que a máquina pega", "Tipologia → sinal automático → onde o humano entra.")
    table(s, MARGIN, 2.8, [3.05, 4.35, PW - 2 * MARGIN - 3.05 - 4.35],
          ["TIPOLOGIA", "SINAL DETECTÁVEL AUTOMÁTICO", "LIMITE — PRECISA DE HUMANO"],
          [["Marcação inflada", "preço desvia da referência de mercado", "laudo nível 3 depende de avaliador"],
           ["Ativo sem lastro (FIDC)", "posição não bate; NF-e/SPED; concentração", "existência do crédito = auditoria física"],
           ["Partes relacionadas", "grafo liga a contraparte ao gestor", "vínculo informal, fora do dado"],
           ["Lavagem / laranjas", "aplica-resgata sem lógica; rede de vínculos", "conluio; documento cadastral forjado"],
           ["Desenquadramento / liquidez", "limite ou prazo violado; descasamento", "determinístico — pega bem"],
           ["Cota / retorno anômalo", "série fora de N desvios; retorno liso", "quem faz “as contas fecharem”"]],
          rh=0.44, csize=9, dsize=10.5)
    txt(s, MARGIN, 6.32, PW - 2 * MARGIN, 0.5,
        "O motor cobre o determinístico e a priorização por risco; a verificação física de lastro, o julgamento de "
        "conluio e o laudo de ilíquido ficam com o humano — desenho híbrido, com trilha de auditoria.",
        10, font=SANS, color=MUTED, italic=True, spacing=1.15)
    source(s, "guia_controle_fraudes_riscos.md; Res. CVM 50/2021 (art. 20); literatura de AML (FATF; dataset Elliptic).")

def s_ax_casos():
    s = slide(PAPER)
    head(s, "Anexo · vigilância", "Casos reais", "O padrão se repete — e o administrador também é punido.")
    table(s, MARGIN, 2.8, [2.95, 0.7, 4.2, PW - 2 * MARGIN - 2.95 - 0.7 - 4.2],
          ["CASO", "ANO", "MECANISMO", "SINAL NOS DADOS"],
          [["Silverado / FIDCs", "2024", "recebíveis fictícios; cedentes ligados à gestora", "lastro; grafo cedente↔gestor"],
           ["Cruzeiro do Sul", "2012", "~320 mil consignados falsos cedidos a FIDCs", "existência de CPF/contrato"],
           ["Banco Master*", "2023–25", "triangula caixa; ativo de R$ 850 mi vira R$ 10 bi", "fluxo circular; salto de marcação"],
           ["FIP LSH", "2024", "laudo do único ativo superfaturado", "salto de marcação sem evento"],
           ["Front-running BB Asset", "2025", "ordens dos fundos vazadas a familiares", "contraparte e timing recorrentes"],
           ["GAS · faraó dos bitcoins", "2021", "Ponzi cripto (10%/mês); ~R$ 17 bi em 12 meses", "rendimento impossível"]],
          rh=0.44, csize=9, dsize=10.5)
    txt(s, MARGIN, 6.32, PW - 2 * MARGIN, 0.5,
        "* Banco Master: caso em apuração (2025–26); mecanismo conforme autoridades e imprensa. Os demais têm decisão da CVM ou do Bacen.",
        9.5, font=SANS, color=MUTED, italic=True, spacing=1.15)
    source(s, "CVM (PAS Silverado 2024; FIP LSH; BB Asset 2025); Bacen; imprensa — ver anexo de Fontes.")

def s_ax_pld():
    s = slide(PAPER)
    head(s, "Anexo · regulatório", "Obrigações de PLD/FT", "Res. CVM 50/2021 — o administrador é sujeito obrigado.")
    itens = [
        ("Política e risco", "política de PLD/FT documentada + abordagem baseada em risco (baixo/médio/alto)."),
        ("Conheça o cliente", "KYC e beneficiário final até a pessoa natural; cadastro atualizado em ≤ 5 anos."),
        ("Monitorar atípicas", "rastrear o rol de operações suspeitas do art. 20, independentemente do valor."),
        ("Comunicar ao COAF", "em até 24 h da conclusão da análise, sob sigilo — sem avisar o comunicado."),
        ("Declaração negativa", "se nada houver a comunicar, declarar à CVM uma vez por ano (até abril)."),
        ("Guardar a prova", "documentação e conclusões à disposição da CVM por no mínimo 5 anos."),
    ]
    cw = (PW - 2 * MARGIN - 0.5) / 2
    for i, (t, d) in enumerate(itens):
        col = i % 2; row = i // 2
        x = MARGIN + col * (cw + 0.5); y = 2.95 + row * 1.02
        oval(s, x + 0.06, y + 0.12, 0.05, fill=GOLD)
        txt(s, x + 0.26, y, cw - 0.26, 0.32, t, 13, font=SERIF, color=INK, bold=True)
        txt(s, x + 0.26, y + 0.36, cw - 0.26, 0.6, d, 10.5, font=SANS, color=MUTED, spacing=1.2)
    rect(s, MARGIN, 5.98, PW - 2 * MARGIN, 0.5, fill=SOFT, line=GOLD, radius=0.05, shadow=False)
    txt(s, MARGIN + 0.25, 6.07, PW - 2 * MARGIN - 0.5, 0.36,
        "A plataforma embute KYC, monitoramento e trilha — o PLD deixa de ser custo de equipe e vira função de software.",
        10.5, font=SANS, color=INK, italic=True, spacing=1.1)
    source(s, "Res. CVM 50/2021 (arts. 4–6, 11–17, 20, 22, 23, 26); Guia ANBIMA PLD/FTP 2025.")

def s_ax_banco():
    s = slide(PAPER)
    head(s, "Anexo · parceria", "O que o banco faz", "Seis passos — e nós ajudamos em todos.")
    reqs = [
        ("Habilitar-se administrador", "protocolo na CVM (SSM/ANBIMA), ~60 dias — com os documentos que entregamos prontos."),
        ("Designar diretores", "responsável pela administração fiduciária + controles internos (pessoas reais da casa)."),
        ("Validar e assinar", "os manuais, o formulário de referência e o website que preparamos."),
        ("Aderir ao Código ANBIMA", "no mesmo protocolo da habilitação."),
        ("Decidir a custódia", "contratar um terceiro no início ou licenciar-se (Res. 32) quando houver escala."),
        ("Supervisionar", "compliance com acesso total, relatórios periódicos e direito de auditar."),
    ]
    x0 = MARGIN; cw = (PW - 2 * MARGIN - 2 * 0.4) / 3
    for i, (t, d) in enumerate(reqs):
        col = i % 3; row = i // 3
        x = x0 + col * (cw + 0.4); y = 2.75 + row * 1.75
        card(s, x, y, cw, 1.55, radius=0.04)
        txt(s, x + 0.26, y + 0.2, 0.7, 0.4, f"{i+1:02d}", 13, font=MONO, color=GOLD, bold=True)
        txt(s, x + 0.9, y + 0.19, cw - 1.1, 0.4, t, 12.5, font=SERIF, color=INK, bold=True, spacing=1.0)
        txt(s, x + 0.26, y + 0.66, cw - 0.5, 0.8, d, 10, font=SANS, color=MUTED, spacing=1.22)
    txt(s, MARGIN, 6.35, PW - 2 * MARGIN, 0.4,
        "Custos ao banco: registro CVM ~R$ 2.380 (uma vez) e fiscalização ~R$ 9.519/ano (administrador); "
        "custódia ~R$ 38.077/ano só se internalizar.", 9.5, font=SANS, color=MUTED, italic=True, spacing=1.1)
    source(s, "guia_burocratico_regulatorio.md; guia_credenciamento_banco.md.")

def s_ax_fontes():
    s = slide(PAPER)
    head(s, "Anexo · metodologia", "Fontes", "De onde vêm os números.")
    groups = [
        ("Modelo e economia", "Documentos internos do projeto (auditados): plano_dtvm_fundos_pequenos.md, guia_potencial_financeiro.md, planilha_custos.md — taxa 0,08%+R$100, viabilidade do cotista, custos."),
        ("Regulatório", "CVM — Resoluções 21 (administração), 32 (custódia), 50 (PLD) e 175 (estrutura); Lei 14.317 (taxa CVM). Pacotes em custodiante/ e administrador_fiduciario/."),
        ("Mercado e nicho", "ANBIMA (~31 mil fundos, 976 gestoras, 51% <R$400mi); ANCORD (27.515 assessores); InfoMoney (custo de fundo pequeno R$ 50–70 mil/ano)."),
        ("Concorrência e juros", "Sites e imprensa sobre Kanastra, QI Tech, Vórtx, BRLTrust, Singulare (foco FIDC); BCB/Copom (Selic 14,25%, jun/2026)."),
    ]
    y = 2.7
    dw = PW - 2 * MARGIN - 3.5 - 0.35
    for t, d in groups:
        txt(s, MARGIN, y, 3.3, 0.5, t, 13, font=SERIF, color=INK, bold=True, spacing=1.05)
        txt(s, MARGIN + 3.5, y, dw, 0.9, d, 10.5, font=SANS, color=MUTED, spacing=1.24)
        y += 0.9
        hline(s, MARGIN, y - 0.14, PW - 2 * MARGIN, color=LINE_DK, weight=0.75)
    rich(s, MARGIN, y + 0.06, PW - 2 * MARGIN - 0.35, 0.8, [
        ([("Honestidade. ", {'size': 10.5, 'font': SANS, 'color': INK, 'bold': True}),
          ("Números de mercado têm data-base citada; a receita é aritmética sobre o PL (não projeção de PL); "
           "a viabilidade do cotista assume CDI+1 e é conservadora com a Selic atual. O protótipo usa dados "
           "simulados e não substitui a licença nem a parceria proposta.",
           {'size': 10.5, 'font': SANS, 'color': MUTED})], {'line': 1.26})])

# ---------------------------------------------------------------- ordem do deck
for fn in [s_cover, s_problema, s_lacuna, s_oportunidade, s_tese, s_solucao, s_deploy,
           s_captacao, s_subclasses,
           s_preco, s_barateamento, s_cotista, s_modelo, s_conta_banco,
           s_vigilancia_porque, s_vigilancia_como, s_risco_banco, s_piloto,
           s_regulatorio, s_roadmap, s_riscos, s_pedido, s_fecho,
           s_anexo_divider, s_ax_cotista, s_ax_custo_fundo, s_ax_receita_pl, s_ax_estrutura,
           s_ax_taxacvm, s_ax_concorrencia, s_ax_mercado, s_ax_precos, s_ax_papeis,
           s_ax_monitoramento, s_ax_casos, s_ax_pld, s_ax_banco, s_ax_fontes]:
    fn()

out = r"C:\Users\vinic\OneDrive\Desktop\dtvm\apresentacao\Argus_Apresentacao.pptx"
prs.save(out)
print("OK ->", out, "| slides:", len(prs.slides._sldIdLst))
