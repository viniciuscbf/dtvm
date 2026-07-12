# -*- coding: utf-8 -*-
"""
Argus — One-pagers de abordagem, personalizados por tipo de instituição.
v4: denso em texto, centrado na TESE ECONÔMICA — por que fundo pequeno é inviável
hoje (~R$ 20 mi de PL mínimo), de onde vem o barateamento (custódia própria +
IA na operação + IA no PLD) e a taxa mais competitiva do mercado (viável a R$ 1 mi).
Gera 1 PPTX A4-retrato por tipo (Banco / CTVM / Financeira), no visual da marca.
Depois exporte para PDF (PowerPoint COM) — ver README/abordagem.

Rodar: python onepagers.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

EMU_IN = 914400
NAVY   = RGBColor(0x6a, 0x50, 0xac); NAVY_DK = RGBColor(0x17, 0x12, 0x3a)  # roxo / roxo escuro (marca)
TEAL   = RGBColor(0x4a, 0x6d, 0xc0); GOLD = RGBColor(0x1f, 0xc0, 0xef)     # azul / ciano (marca)
INK    = RGBColor(0x2b, 0x32, 0x32); MUTED = RGBColor(0x5b, 0x6b, 0x7f)
FAINT  = RGBColor(0x90, 0x9d, 0xac); LINE = RGBColor(0xe7, 0xe8, 0xf2)
PAPER  = RGBColor(0xff, 0xff, 0xff); SOFT = RGBColor(0xf4, 0xf5, 0xfb)
D_TXT  = RGBColor(0xf3, 0xf6, 0xfa); D_MUT = RGBColor(0xa8, 0xb6, 0xc6); D_FAINT = RGBColor(0x6d, 0x80, 0x96)
D_LINE = RGBColor(0x2e, 0x26, 0x52)
SERIF, SANS, MONO = "Georgia", "Calibri", "Consolas"
PW, PH, M = 8.27, 11.69, 0.7  # A4 retrato

BASE = os.path.dirname(os.path.abspath(__file__))

def _run(run, t, size, font, color, bold, italic, char):
    run.text = t; f = run.font
    f.name = font; f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color
    if char is not None:
        run._r.get_or_add_rPr().set('spc', str(int(char * 100)))

def txt(s, x, y, w, h, text, size, *, font=SANS, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0, char=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(0); p.space_before = Pt(0)
        _run(p.add_run(), line, size, font, color, bold, italic, char)
    return tb

def rich(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, (runs, o) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = o.get('align', PP_ALIGN.LEFT); p.line_spacing = o.get('line', 1.0)
        p.space_before = Pt(o.get('sb', 0)); p.space_after = Pt(o.get('sa', 0))
        for (t, oo) in runs:
            _run(p.add_run(), t, oo.get('size', 12), oo.get('font', SANS), oo.get('color', INK),
                 oo.get('bold', False), oo.get('italic', False), oo.get('char'))
    return tb

def rect(s, x, y, w, h, fill=None, line=None, lw=0.75, radius=None):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False; return shp

def hline(s, x, y, w, color=LINE, weight=0.75):
    from pptx.enum.shapes import MSO_CONNECTOR
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(weight); ln.shadow.inherit = False; return ln

def oval(s, cx, cy, r, fill=None, line=None, lw=1.0):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False; return shp

MARK = r"C:\Users\vinic\OneDrive\Desktop\dtvm\brand\favicon_argus.png"

def kicker(s, x, y, text):
    hline(s, x, y + 0.08, 0.24, color=GOLD, weight=1.6)
    txt(s, x + 0.34, y, 6.6, 0.25, text.upper(), 8.5, font=MONO, color=GOLD, char=1.8)

def stat_cards(s, x, y, w, itens):
    n = len(itens); gap = 0.16; cw = (w - gap * (n - 1)) / n
    for i, (v, r) in enumerate(itens):
        xx = x + i * (cw + gap)
        rect(s, xx, y, cw, 0.68, fill=SOFT, line=LINE, radius=0.09)
        txt(s, xx + 0.14, y + 0.08, cw - 0.28, 0.3, v, 12.5, font=SERIF, color=NAVY_DK, bold=True)
        txt(s, xx + 0.14, y + 0.4, cw - 0.28, 0.26, r, 7.3, font=SANS, color=MUTED, spacing=1.0)

def bloco(s, x, y, w, titulo, texto):
    oval(s, x + 0.05, y + 0.09, 0.04, fill=GOLD)
    txt(s, x + 0.2, y, w - 0.2, 0.2, titulo, 9.8, font=SANS, color=NAVY_DK, bold=True)
    txt(s, x + 0.2, y + 0.2, w - 0.2, 0.5, texto, 9.0, font=SANS, color=INK, spacing=1.08)

def onepager(cfg):
    prs = Presentation(); prs.slide_width = Emu(int(PW*EMU_IN)); prs.slide_height = Emu(int(PH*EMU_IN))
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = PAPER
    W = PW - 2*M

    # faixa superior escura — logo oficial completa (versão de texto branco p/ fundo escuro)
    rect(s, 0, 0, PW, 2.42, fill=NAVY_DK)
    LOGO_BRANCA = r"C:\Users\vinic\OneDrive\Desktop\dtvm\brand\logo_argus_branca.png"
    s.shapes.add_picture(LOGO_BRANCA, Inches(M), Inches(0.42), height=Inches(0.54))
    txt(s, M + 0.02, 1.05, 6, 0.3, "ADMINISTRAÇÃO FIDUCIÁRIA E CUSTÓDIA DE FUNDOS", 7.5, font=MONO, color=D_MUT, char=1.5)
    rich(s, M, 1.36, W, 0.8, [([(cfg['titulo'], {'size': 19, 'font': SERIF, 'color': D_TXT})], {'line': 1.06})])
    txt(s, M, 2.1, W, 0.28, cfg['sub'], 9.5, font=SANS, color=GOLD, italic=True)

    # ---- A ideia + a promessa ----
    y = 2.58
    kicker(s, M, y, "A ideia, direto ao ponto")
    txt(s, M, y + 0.28, W, 0.74,
        "Estamos estruturando uma operação de DTVM para o vazio que os grandes administradores criaram: fundos de "
        "investimento de menor patrimônio. Com tecnologia própria, constituímos e administramos fundos de forma "
        "rápida e barata — com o objetivo declarado de operar a taxa de administração fiduciária MAIS COMPETITIVA "
        "do mercado brasileiro, sem abrir mão de nenhum controle fiduciário. Só 1.277 dos ~31 mil fundos do país "
        "passam de R$ 1 bilhão: a cauda é enorme e está sem dono.", 9.6, font=SANS, color=INK, spacing=1.14)

    # ---- Por que é inviável hoje ----
    y = 3.66
    kicker(s, M, y, "Por que fundo pequeno hoje é inviável")
    txt(s, M, y + 0.28, W, 0.9,
        "Hoje, constituir um fundo só fecha a conta a partir de ~R$ 20 milhões de patrimônio: os pisos das "
        "administradoras tradicionais (R$ 12–25 mil/mês somando administração, custódia e controladoria) impõem "
        "R$ 150–300 mil de custo fixo por ano. Num fundo de R$ 5 milhões, isso consome 3–6% do patrimônio ao ano "
        "antes de qualquer resultado — inviável por definição. O gestor emergente fica preso a clube de investimento "
        "ou carteira administrada. O nosso barateamento é grande o suficiente para inverter essa conta: fundo "
        "viável a partir de R$ 1 milhão de PL.", 9.6, font=SANS, color=INK, spacing=1.14)
    stat_cards(s, M, y + 1.28, W, [
        ("R$ 150–300 mil/ano", "custo fixo típico do fundo hoje"),
        ("≈ R$ 20 milhões", "PL mínimo p/ fechar a conta hoje"),
        ("R$ 1 milhão", "PL viável na estrutura Argus")])

    # ---- De onde vem o barateamento ----
    y = 5.86
    kicker(s, M, y, "De onde vem o barateamento — três alavancas")
    by = y + 0.3
    bloco(s, M, by, W, cfg['alav1_titulo'], cfg['alav1'])
    bloco(s, M, by + 0.78, W, "Operação automatizada com IA — custo marginal por fundo perto de zero.",
          "Cota diária, conciliação, enquadramento, tributação, reporte CVM/ANBIMA e a geração dos documentos de "
          "constituição rodam automatizados; a IA assume a conferência e as exceções que hoje ocupam equipes "
          "inteiras. Administrar o fundo nº 100 custa quase o mesmo que o nº 10.")
    bloco(s, M, by + 1.56, W, "Verificação antilavagem com IA — vigilância diária de 100% dos fundos.",
          "O motor roda regras de preço fora de mercado, partes relacionadas, lastro em custódia e movimentação "
          "atípica todos os dias, com evidência e trilha. O compliance que nas casas tradicionais é custo fixo de "
          "gente vira software — mais barato e MAIS rigoroso que o padrão de mercado.")

    # ---- Banner taxa ----
    y = 8.5
    rect(s, M, y, W, 0.86, fill=NAVY, radius=0.05)
    txt(s, M + 0.3, y + 0.13, 4.4, 0.36, "0,08% a.a. · piso R$ 100/mês", 15.5, font=SERIF, color=D_TXT, bold=True)
    txt(s, M + 0.3, y + 0.53, 4.6, 0.28, "A taxa mais competitiva do mercado — 10–40× abaixo dos pisos praticados.", 8.3, font=SANS, color=D_MUT)
    txt(s, PW - M - 2.75, y + 0.15, 2.45, 0.6, "fundo de R$ 1 milhão:\n≈ R$ 1.200/ano — contra\nR$ 150–300 mil hoje", 8.3,
        font=MONO, color=RGBColor(0xcf, 0xee, 0xfb), char=0.2, align=PP_ALIGN.RIGHT, spacing=1.18)

    # ---- A parceria (por segmento) ----
    y = 9.5
    kicker(s, M, y, cfg['parceria_titulo'])
    txt(s, M, y + 0.27, W, 0.8, cfg['parceria'], 9.2, font=SANS, color=INK, spacing=1.1)

    # ---- convite humano ----
    y = 10.64
    rect(s, M, y, W, 0.44, fill=SOFT, line=GOLD, lw=1.0, radius=0.14)
    txt(s, M + 0.26, y + 0.08, W - 0.52, 0.3, cfg['cta'], 9.3, font=SERIF, color=INK, italic=True, spacing=1.08)

    # rodapé
    hline(s, M, PH - 0.46, W, color=LINE)
    txt(s, M, PH - 0.35, 4.6, 0.25, "vinicius.fernandes@argusdtvm.com.br  ·  argusdtvm.com.br", 8.5, font=MONO, color=NAVY, char=0.2)
    txt(s, PW - M - 2.9, PH - 0.35, 2.9, 0.25, "CONFIDENCIAL · DADOS SIMULADOS",
        7, font=MONO, color=FAINT, char=0.5, align=PP_ALIGN.RIGHT)

    out = os.path.join(BASE, cfg['arquivo']); prs.save(out); return out

CFGS = [
    {"arquivo": "Argus_OnePager_Banco.pptx",
     "titulo": "Fundos viáveis a partir de R$ 1 milhão —\ne o papel do seu banco nisso.",
     "sub": "Onde nasce o custo que mata o fundo pequeno, como ele desaparece, e o que propomos ao banco.",
     "alav1_titulo": "Custódia dentro de casa — o banco como agente custodiante.",
     "alav1": ("Em vez de pagar custódia de terceiro — que sozinha já carrega piso de mercado —, a parceria credencia a "
               "custódia no próprio banco (Res. CVM 32): o custo sai da conta do fundo e vira receita interna da "
               "instituição. É a verticalização que os grandes praticam, aplicada pela primeira vez ao fundo pequeno."),
     "parceria_titulo": "O que propomos ao seu banco",
     "parceria": ("Administração fiduciária e custódia são privativas de instituição autorizada pelo BCB — e é aí que "
                  "entra o banco: licença, estrutura jurídica e o compliance que já exerce. Nós entramos com tecnologia, "
                  "equipe, operação e todos os custos (não buscamos investimento); a contrapartida é participação nos "
                  "lucros, sem o banco absorver despesa. Caminho: carta de intenções → credenciamento na CVM (dossiê "
                  "pronto, prazo legal de 60 dias) → produção sob a marca do banco → custódia própria (Res. 32). "
                  "Piloto no ar: argusdtvm.com.br/piloto — o tour anexo percorre cada tela."),
     "cta": "Queremos abrir essa conta com os números do seu banco — 20 minutos com quem construiu a plataforma. É só responder o e-mail."},

    {"arquivo": "Argus_OnePager_CTVM.pptx",
     "titulo": "Fundos viáveis a partir de R$ 1 milhão —\ne o papel da sua corretora nisso.",
     "sub": "Onde nasce o custo que mata o fundo pequeno, como ele desaparece, e o que propomos à corretora.",
     "alav1_titulo": "Custódia dentro de casa — a corretora como agente custodiante.",
     "alav1": ("Em vez de pagar custódia de terceiro — que sozinha já carrega piso de mercado —, a parceria credencia a "
               "custódia na própria corretora (a CTVM é elegível às duas licenças, Res. CVM 21 e 32): o custo sai da "
               "conta do fundo e vira receita interna. Verticalização dos grandes, aplicada ao fundo pequeno — e a "
               "corretagem dos fundos da casa fica em casa também."),
     "parceria_titulo": "O que propomos à sua corretora",
     "parceria": ("Administração fiduciária é privativa de instituição autorizada — e a licença de CTVM já habilita a "
                  "corretora: nada novo a constituir, só o credenciamento na CVM. Vocês entram com a licença, o "
                  "compliance e a base de gestores da casa; nós, com tecnologia, equipe, operação e todos os custos "
                  "(não buscamos investimento) — contrapartida em participação nos lucros, sem absorver despesa. "
                  "Caminho: carta de intenções → credenciamento (dossiê pronto, ~60 dias) → produção sob a marca da "
                  "corretora → custódia própria (Res. 32). Piloto no ar: argusdtvm.com.br/piloto — tour anexo."),
     "cta": "Queremos abrir essa conta com os números da sua corretora — 20 minutos com quem construiu a plataforma. É só responder o e-mail."},

    {"arquivo": "Argus_OnePager_Financeira.pptx",
     "titulo": "Fundos viáveis a partir de R$ 1 milhão —\ne a janela da Res. 5.237 para a financeira.",
     "sub": "Onde nasce o custo que mata o fundo pequeno, como ele desaparece, e o que propomos à financeira.",
     "alav1_titulo": "Custódia estruturada para não pesar no fundo.",
     "alav1": ("A custódia nasce contratada em bloco com custodiante autorizado — negociada em escala e operada pela "
               "nossa plataforma, o custo por fundo cai a uma fração do avulso. E a arquitetura já nasce pronta para "
               "internalizar a custódia num parceiro habilitado quando a escala justificar, eliminando o repasse por "
               "completo."),
     "parceria_titulo": "O que propomos à sua financeira",
     "parceria": ("Administração fiduciária é privativa de instituição autorizada — e desde 1º/9/2025 a Res. CMN 5.237 "
                  "(art. 6º, p.u., IV) habilitou as financeiras: a sua licença passou a servir exatamente para isso, e o "
                  "mercado ainda não acordou. Vocês entram com a instituição e o compliance; nós, com tecnologia, equipe, "
                  "operação e todos os custos (não buscamos investimento) — contrapartida em participação nos lucros, sem "
                  "absorver despesa. Caminho: carta de intenções → registro na CVM pela 5.237 (dossiê pronto, ~60 dias) → "
                  "produção sob a marca da financeira. Piloto no ar: argusdtvm.com.br/piloto — tour anexo."),
     "cta": "Queremos mostrar a Res. 5.237 aplicada aos números da sua financeira — 20 minutos com quem construiu a plataforma. É só responder o e-mail."},
]

if __name__ == "__main__":
    for c in CFGS:
        print("gerado:", onepager(c))
